import os
import fitz # PyMuPDF for PDF
from docx import Document
import pandas as pd
from pptx import Presentation
from transformers import pipeline
import torch
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
import shutil
import json
from collections import Counter

# --- Model Selection Function ---
def select_summarization_model(last_model=None):
    """
    Prompts the user to select a summarization model from a list of options.
    Optionally offers the last-used model.
    """
    if last_model:
        choice = input(f"Use last summarization model '{last_model}'? (y/n) [default: y]: ").strip().lower()
        if choice in ['y', 'yes', '']:
            return last_model
    
    print("\nPlease select a summarization model to use:")
    print("--- Models Sorted by Size (Largest to Smallest) ---")
    print("1: (2440 MB) facebook/mbart-large-cc25 - Compact multilingual")
    print("2: (2330 MB) csebuetnlp/mT5_multilingual_XLSum - Summarizes in 44 languages")
    print("3: (2280 MB) google/pegasus-xsum - Headline-style summaries")
    print("4: (1630 MB) facebook/bart-large-cnn - Larger, higher quality model")
    print("5: (1630 MB) philschmid/bart-large-cnn-samsum - Specialized for conversations/dialogue")
    print("6: (990 MB) google/flan-t5-base - Excellent for complex text & instructions")
    print("7: (892 MB) mrm8488/t5-base-finetuned-summarize-news - Specialized for news articles")
    print("8: (496 MB) sshleifer/distilbart-cnn-12-6 - Faster, distilled model for speed")
    print("9: (496 MB) ainize/kobart-news - Specialized for Korean news")
    print("10: (242 MB) t5-small - Balanced speed and quality (Default)")
    print("11: (242 MB) moussaKam/t5-small-finetuned-xsum - General purpose, reliable")
    print("12: (242 MB) DivyanshuSheth/T5-Seq2Seq-Final - General purpose (T5-small)")
    print("13: (242 MB) Shahm/t5-small-german - Specialized for German text")
    print("14: (242 MB) efederici/text2tags - For generating tags/keywords")
    print("15: (242 MB) Falconsai/medical_summarization - High-performing for medical text")
    print("16: (134 MB) Aryan0310/bert-mini2bert-mini-finetuned-cnn_daily_mail-summarization-finetuned-xsum - Tuned for news/headlines")
    print("17: (125 MB) google/t5-efficient-mini - Most resource-efficient")
    print("18: (125 MB) google/t5-efficient-mini-xsum - Tuned for headline summaries")
    print("19: (125 MB) google/t5-efficient-mini-cnndm - Tuned for news")
    print("20: (125 MB) google/t5-efficient-mini-samsum - Tuned for dialogue")
    
    models = {
        "1": "facebook/mbart-large-cc25",
        "2": "csebuetnlp/mT5_multilingual_XLSum",
        "3": "google/pegasus-xsum",
        "4": "facebook/bart-large-cnn",
        "5": "philschmid/bart-large-cnn-samsum",
        "6": "google/flan-t5-base",
        "7": "mrm8488/t5-base-finetuned-summarize-news",
        "8": "sshleifer/distilbart-cnn-12-6",
        "9": "ainize/kobart-news",
        "10": "t5-small",
        "11": "moussaKam/t5-small-finetuned-xsum",
        "12": "DivyanshuSheth/T5-Seq2Seq-Final",
        "13": "Shahm/t5-small-german",
        "14": "efederici/text2tags",
        "15": "Falconsai/medical_summarization",
        "16": "Aryan0310/bert-mini2bert-mini-finetuned-cnn_daily_mail-summarization-finetuned-xsum",
        "17": "google/t5-efficient-mini",
        "18": "google/t5-efficient-mini-xsum",
        "19": "google/t5-efficient-mini-cnndm",
        "20": "google/t5-efficient-mini-samsum"
    }
    
    while True:
        choice = input("Enter your choice (1-20): ").strip()
        if choice in models:
            return models[choice]
        elif choice == "": # Default option
            return models["10"]
        print("Invalid choice. Please enter a number from 1 to 20.")

# --- NEW FUNCTION: Classifier Model Selection ---
def select_classifier_model(last_model=None):
    """
    Prompts the user to select a classifier model from a list of options.
    Optionally offers the last-used model.
    """
    if last_model:
        choice = input(f"Use last classifier model '{last_model}'? (y/n) [default: y]: ").strip().lower()
        if choice in ['y', 'yes', '']:
            return last_model
    
    print("\nPlease select a classifier model:")
    print("--- Models Sorted by Size (Largest to Smallest) ---")
    print("1: (440 MB) openai-community/roberta-large-openai-detector - Detects AI-generated text")
    print("2: (440 MB) FacebookAI/roberta-large-mnli - Multi-genre natural language inference (Good for scientific fields)")
    print("3: (440 MB) amannor/bert-base-uncased-sdg-classifier - General purpose BERT model (Good for scientific fields)")
    print("4: (377 MB) HugoGiddins/multi-tag-classifier-debertav3-small-2 - Powerful multi-tag classifier (Good for scientific fields)")
    print("5: (268 MB) nori3tsu/classify-reservation-intent - Specialized for user intent")
    print("6: (260 MB) openai-community/roberta-base-openai-detector - Detects AI-generated text")
    print("7: (260 MB) bvanaken/clinical-assertion-negation-bert - Clinical assertion and negation detection (Good for scientific fields)")
    print("8: (260 MB) Falconsai/intent_classification - General user intent classification")
    print("9: (260 MB) TheBritishLibrary/bl-books-genre - Genre classification for books (Good for scientific fields)")
    print("10: (260 MB) aditeyabaral/finetuned-iitp_pdt_review-additionalpretrained-bert-base-cased - Review classification")
    print("11: (260 MB) aloxatel/bert-base-mnli - Multi-genre natural language inference")
    print("12: (260 MB) blackbird/bert-base-uncased-MNLI-v1 - Multi-genre natural language inference")
    print("13: (260 MB) bgoel4132/tweet-disaster-classifier - Disaster tweet classification")
    print("14: (20 MB) bioformers/bioformer-8L-qnli - Question-and-answer model for biomedical text (Good for scientific fields)")
    print("15: (16 MB) bioformers/bioformer-8L-mnli - Multi-genre natural language inference for biomedical text (Good for scientific fields)")
    print("16: (3 MB) celential/erc - General purpose conversational classification")

    models = {
        "1": "openai-community/roberta-large-openai-detector",
        "2": "FacebookAI/roberta-large-mnli",
        "3": "amannor/bert-base-uncased-sdg-classifier",
        "4": "HugoGiddins/multi-tag-classifier-debertav3-small-2",
        "5": "nori3tsu/classify-reservation-intent",
        "6": "openai-community/roberta-base-openai-detector",
        "7": "bvanaken/clinical-assertion-negation-bert",
        "8": "Falconsai/intent_classification",
        "9": "TheBritishLibrary/bl-books-genre",
        "10": "aditeyabaral/finetuned-iitp_pdt_review-additionalpretrained-bert-base-cased",
        "11": "aloxatel/bert-base-mnli",
        "12": "blackbird/bert-base-uncased-MNLI-v1",
        "13": "bgoel4132/tweet-disaster-classifier",
        "14": "bioformers/bioformer-8L-qnli",
        "15": "bioformers/bioformer-8L-mnli",
        "16": "celential/erc"
    }
    
    while True:
        choice = input("Enter your choice (1-16): ").strip()
        if choice in models:
            return models[choice]
        elif choice == "": # Default option
            return models["1"]
        print("Invalid choice. Please enter a number from 1 to 16.")

# --- NEW FUNCTION: Get Summary Lengths ---
def get_summary_lengths():
    """
    Prompts the user to set the min and max summary lengths with validation.
    """
    min_len, max_len = 30, 70 # Default values
    
    while True:
        try:
            min_input = input(f"Enter MIN summary length [default: {min_len}]: ").strip()
            if min_input == "":
                min_len_to_use = min_len
            else:
                min_len_to_use = int(min_input)
            
            max_input = input(f"Enter MAX summary length [default: {max_len}]: ").strip()
            if max_input == "":
                max_len_to_use = max_len
            else:
                max_len_to_use = int(max_input)
            
            # Ensure max_length is always greater than min_length
            if max_len_to_use < min_len_to_use:
                print(f"Error: Max length ({max_len_to_use}) must be greater than min length ({min_len_to_use}). Automatically adjusting max length to {min_len_to_use + 50}.")
                max_len_to_use = min_len_to_use + 50

            break
        except ValueError:
            print("Invalid input. Please enter a valid number.")
            
    print(f"Summary lengths set to MIN: {min_len_to_use}, MAX: {max_len_to_use}")
    return min_len_to_use, max_len_to_use

# --- NEW FUNCTION: Get and remember classifier threshold for each model ---
def get_classifier_threshold(model_name):
    """
    Asks the user for a classifier threshold, remembering the value for each model.
    """
    settings_file = "model_thresholds.json"
    thresholds = {}
    
    # Load existing thresholds from file
    try:
        if os.path.exists(settings_file):
            with open(settings_file, 'r', encoding='utf-8') as f:
                thresholds = json.load(f)
    except Exception as e:
        print(f"Warning: Could not load threshold settings file: {e}")

    # Get the saved threshold for the current model, or use 0.1 as a default
    default_threshold = thresholds.get(model_name, 0.1)

    while True:
        prompt = f"\nEnter classifier threshold for '{model_name}' (0.0-1.0) [default: {default_threshold}]: "
        user_input = input(prompt).strip()

        if user_input == "":
            chosen_threshold = default_threshold
            break

        try:
            chosen_threshold = float(user_input)
            if 0.0 <= chosen_threshold <= 1.0:
                break
            else:
                print("Error: Please enter a number between 0.0 and 1.0.")
        except ValueError:
            print("Error: Invalid input. Please enter a valid number.")

    # Save the chosen threshold for this model
    thresholds[model_name] = chosen_threshold
    try:
        with open(settings_file, 'w', encoding='utf-8') as f:
            json.dump(thresholds, f, indent=4)
        print(f"Threshold for this model set to {chosen_threshold} and saved.")
    except Exception as e:
        print(f"Warning: Could not save threshold settings: {e}")
        
    return chosen_threshold

# --- NEW FUNCTION: Save Last Used Models ---
def save_last_models(summarizer_model_name, classifier_model_name):
    """
    Saves the last used model names to a file.
    """
    settings_file = "last_models.json"
    data = {
        "summarizer_model": summarizer_model_name,
        "classifier_model": classifier_model_name
    }
    try:
        with open(settings_file, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4)
    except Exception as e:
        print(f"Warning: Could not save last model settings: {e}")

# --- NEW FUNCTION: Load Last Used Models ---
def load_last_models():
    """
    Loads the last used model names from a file.
    """
    settings_file = "last_models.json"
    try:
        if os.path.exists(settings_file):
            with open(settings_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return data.get("summarizer_model"), data.get("classifier_model")
    except Exception as e:
        print(f"Warning: Could not load last model settings file: {e}")
    return None, None

# Check for GPU availability and set the device
if torch.cuda.is_available():
    device = "cuda"
    print("Device set to use GPU.")
else:
    device = "cpu"
    print("GPU not found. Device set to use CPU.")

# --- Configuration ---
# The logic for chunking is now handled by the tokenizer, not this variable.
# MAX_CHUNK_LENGTH is now a conceptual size for file reading only.
MAX_CHUNK_LENGTH = 512

# --- NEW FUNCTION: Get Download Directory ---
def get_download_directory():
    """
    Prompts the user for a download directory and creates it if it doesn't exist.
    """
    default_path = os.path.join(os.getcwd(), "models_cache")
    print(f"\nModel files can be large. By default, they are stored in a cache directory.")
    print(f"Current working directory: {os.getcwd()}")
    path_input = input(f"Enter a directory path to store model files [default: '{default_path}']: ").strip()

    if path_input == "":
        cache_dir = default_path
    else:
        cache_dir = path_input

    try:
        os.makedirs(cache_dir, exist_ok=True)
        print(f"Model download directory set to: {cache_dir}")
        return cache_dir
    except Exception as e:
        print(f"Error creating directory: {e}. Falling back to default cache location.")
        return None

# --- Get Token Chunk Size Function ---
def get_token_chunk_size():
    """
    Prompts the user to set the token chunk size for models.
    """
    default_chunk_size = 140
    while True:
        try:
            user_input = input(f"Enter the maximum token chunk size for models [default: {default_chunk_size}]: ").strip()
            if user_input == "":
                return default_chunk_size
            else:
                size = int(user_input)
                if size > 0:
                    return size
                else:
                    print("Error: Please enter a number greater than 0.")
        except ValueError:
            print("Error: Invalid input. Please enter a valid number.")

# --- Load and Edit Categories Function ---
def load_and_edit_categories():
    """
    Loads categories from a file, prompts the user to edit them, and saves the final list.
    """
    category_file = "categories.json"
    default_categories = [
        "Health",
        "History",
        "Weather",
        "Technology",
        "Finance"
    ]
    categories = default_categories
    
    try:
        if os.path.exists(category_file):
            with open(category_file, 'r', encoding='utf-8') as f:
                categories = json.load(f)
            print("Step 1.2: Loaded categories from previous session.")
    except Exception as e:
        print(f"Error loading categories file: {e}. Using default categories.")

    # Display the current categories list here, before the 'edit' prompt
    print("\n--- Current Categories ---")
    for i, cat in enumerate(categories):
        print(f"[{i+1}] {cat}")
    print("--------------------------")
    
    edit_choice = input("Do you want to edit these categories? (y/n) [default: no]: ").strip().lower()
    if edit_choice in ['y', 'yes']:
        while True:
            print("\n--- Current Categories ---")
            for i, cat in enumerate(categories):
                print(f"[{i+1}] {cat}")
            print("--------------------------")
            print("\nEditing categories. Enter 'add <new_category>', 'remove <number>', 'edit <number> <new_category>', 'list', or 'done' to finish.")

            command = input("> ").strip()
            if command.lower() == 'done':
                break
            
            parts = command.split()
            if not parts:
                continue

            action = parts[0].lower()
            if action == 'add' and len(parts) >= 2:
                new_cat = " ".join(parts[1:])
                categories.append(new_cat)
                print(f"Added: {new_cat}")
            elif action == 'remove' and len(parts) == 2 and parts[1].isdigit():
                idx = int(parts[1]) - 1
                if 0 <= idx < len(categories):
                    removed_cat = categories.pop(idx)
                    print(f"Removed: {removed_cat}")
                else:
                    print("Invalid category number.")
            elif action == 'edit' and len(parts) >= 3 and parts[1].isdigit():
                idx = int(parts[1]) - 1
                if 0 <= idx < len(categories):
                    new_cat = " ".join(parts[2:])
                    old_cat = categories[idx]
                    categories[idx] = new_cat
                    print(f"Edited: '{old_cat}' to '{new_cat}'")
                else:
                    print("Invalid category number.")
            elif action == 'list':
                pass # The list will be printed by the loop's design
            else:
                print("Invalid command. Please try again.")
        
    # Save the final list of categories
    with open(category_file, 'w', encoding='utf-8') as f:
        json.dump(categories, f, indent=4)
    print("\nStep 1.3: Categories saved for next session.")
        
    return categories

# --- Tokenizer and Chunking Function ---
def chunk_text_by_tokens(text, tokenizer, max_length):
    """
    Splits text into token-based chunks to prevent the model from failing on long sequences.
    """
    chunks = []
    
    # Tokenize the entire text
    tokenized_text = tokenizer.tokenize(text)
    
    # Split the tokenized text into chunks of max_length
    for i in range(0, len(tokenized_text), max_length):
        chunk = tokenized_text[i:i + max_length]
        
        # Detokenize the chunk to get the original text
        detokenized_chunk = tokenizer.convert_tokens_to_string(chunk)
        chunks.append(detokenized_chunk)
    
    return chunks

# --- File Reading Functions ---
def read_pdf(file_path):
    """Extracts all text from a PDF file."""
    text = ""
    try:
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text += page.get_text()
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return text

def read_docx(file_path):
    """Extracts all text from a DOCX file."""
    text = ""
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error reading DOCX: {e}")
    return text

def read_txt(file_path):
    """Reads all text from a plain TXT file with improved encoding handling."""
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            text = file.read()
    except UnicodeDecodeError:
        print("UTF-8 decode failed. Falling back to latin-1 encoding.")
        with open(file_path, "r", encoding="latin-1", errors='replace') as file:
            text = file.read()
    except Exception as e:
        print(f"Error reading TXT: {e}")
    return text

def read_excel(file_path):
    """Extracts all text from an XLSX (Excel) file."""
    text = ""
    try:
        df = pd.read_excel(file_path, sheet_name=None)
        full_text = ""
        for sheet_name in df:
            full_text += df[sheet_name].to_string(index=False) + "\n"
        text = full_text
    except Exception as e:
        print(f"Error reading XLSX: {e}")
    return text

def read_pptx(file_path):
    """Extracts all text from a PPTX (PowerPoint) file."""
    text = ""
    try:
        presentation = Presentation(file_path)
        full_text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text += shape.text + "\n"
        text = full_text
    except Exception as e:
        print(f"Error reading PPTX: {e}")
    return text

# --- Summarization Function ---
def summarize_text(text, summarizer_pipeline, min_len, max_len, chunk_size):
    """
    Generates a summary from the text using the provided summarizer pipeline.
    """
    if not text.strip():
        return ["No text to summarize."]

    # Use the pipeline's tokenizer for chunking to avoid token size issues
    chunks = chunk_text_by_tokens(text, summarizer_pipeline.tokenizer, chunk_size)
    
    summaries = []
    
    print(f"Step 6.1.1: Summarizing {len(chunks)} chunks in a batch...")
    try:
        results = summarizer_pipeline(
            chunks, 
            max_new_tokens=max_len, 
            min_length=min_len, 
            truncation=True,
            batch_size=8
        )
        summaries = [result['summary_text'] for result in results]
        print(f"Step 6.1.2: Summarization for all chunks complete.")
    except Exception as e:
        print(f"Step 6.1.3: Error during batch summarization: {e}")
        for i, chunk in enumerate(chunks):
            try:
                summary = summarizer_pipeline(
                    chunk, max_new_tokens=max_len, min_length=min_len, truncation=True
                )
                summaries.append(summary[0]['summary_text'])
            except Exception as chunk_e:
                print(f"Error on chunk {i+1}: {chunk_e}")
                summaries.append(f"Summary failed for chunk {i+1}.")

    combined_summary = ". ".join(summaries)
    return [s.strip() for s in combined_summary.split(".") if s.strip()]
    
# --- Categorization Function ---
def categorize_summary(summary_text, categories, classifier_pipeline, threshold, chunk_size):
    """Categorizes a summary by chunking it and returning the most common label."""
    if classifier_pipeline is None:
        return "Other"

    print("Step 7.1.1: Sending summary to classifier model...")
    
    try:
        # Use the classifier's tokenizer to create chunks from the summary
        chunks = chunk_text_by_tokens(summary_text, classifier_pipeline.tokenizer, chunk_size)
        
        # Classify each chunk
        chunk_results = []
        for chunk in chunks:
            result = classifier_pipeline(chunk, candidate_labels=categories)
            if result['scores'][0] >= threshold:
                chunk_results.append(result['labels'][0])
        
        # If no chunks were classified above the threshold, return "Other"
        if not chunk_results:
            return "Other"
        
        # Find the most common category among the classified chunks
        most_common_category = Counter(chunk_results).most_common(1)[0][0]
        return most_common_category
    
    except Exception as e:
        print(f"Step 7.1.3: Error during AI categorization: {e}")
        return "Other"

# --- Main Processing Logic ---
def process_files_in_folder(folder_path, scan_subdirectories, categories, start_chunk, end_chunk, file_management_settings, summarizer_pipeline, classifier_pipeline, min_len, max_len, classifier_threshold_to_use, token_chunk_size, color_code):
    """
    Walks a folder, processes supported files, and generates summaries.
    """
    supported_extensions = {
        ".pdf": read_pdf, 
        ".docx": read_docx, 
        ".txt": read_txt, 
        ".xlsx": read_excel, 
        ".pptx": read_pptx
    }
    
    print("\nStep 2: Starting folder scan...")
    file_list = []
    if scan_subdirectories:
        for root, _, files in os.walk(folder_path):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in supported_extensions:
                    file_list.append(os.path.join(root, file))
    else:
        for file in os.listdir(folder_path):
            file_path = os.path.join(folder_path, file)
            if os.path.isfile(file_path):
                ext = os.path.splitext(file)[1].lower()
                if ext in supported_extensions:
                    file_list.append(file_path)

    if not file_list:
        print("Step 3: No supported files found. Exiting.")
        return []

    print(f"Step 4: Found {len(file_list)} supported files. Processing them now...")

    all_summaries = []
    
    # ANSI escape codes for formatting
    BOLD_START = "\033[1m"
    BOLD_END = "\033[0m"
    COLOR_START = color_code
    COLOR_END = "\033[0m"

    for i, file_path in enumerate(file_list):
        print(f"\nStep 5: Processing file {i + 1}/{len(file_list)} - {os.path.basename(file_path)}")
        
        # Read the full text first
        reader = supported_extensions[os.path.splitext(file_path)[1].lower()]
        full_text = reader(file_path)

        if full_text.strip():
            # Get the chunks based on the requested range
            all_chunks = chunk_text_by_tokens(full_text, summarizer_pipeline.tokenizer, token_chunk_size)
            
            # Select the specific chunks the user requested
            selected_chunks = all_chunks[start_chunk-1:end_chunk]

            if not selected_chunks:
                print(f"Warning: No text found in chunks {start_chunk} to {end_chunk}. Skipping file.")
                continue

            print(f"Step 6: Extracting text from chunks {start_chunk} to {end_chunk}...")
            text_to_summarize = " ".join(selected_chunks)

            print("Step 7: Text extracted successfully. Starting summarization...")
            bullet_points = summarize_text(text_to_summarize, summarizer_pipeline, min_len, max_len, token_chunk_size)
            
            summary_text = " ".join(bullet_points)
            print("Step 7.1: Categorizing summary...")
            category = categorize_summary(summary_text, categories, classifier_pipeline, classifier_threshold_to_use, token_chunk_size)

            print("Step 8: Final summary complete.")
            
            print(f"\nFile: {file_path}\nCategory: {category}\nSummary:\n")

            # Bold and colorize the matching category text in the summary for display
            display_summary = summary_text
            for cat in categories:
                if cat.lower() in display_summary.lower():
                    start_index = 0
                    while True:
                        start_index = display_summary.lower().find(cat.lower(), start_index)
                        if start_index == -1:
                            break
                        end_index = start_index + len(cat)
                        original_text = display_summary[start_index:end_index]
                        display_summary = display_summary[:start_index] + COLOR_START + BOLD_START + original_text + COLOR_END + BOLD_END + display_summary[end_index:]
                        start_index += len(COLOR_START) + len(BOLD_START) + len(COLOR_END) + len(BOLD_END) + len(cat)
            
            # Print the formatted summary
            print(display_summary)
            print("\n" + "-"*50)

            summary_text_for_file = ""
            for point in bullet_points:
                summary_text_for_file += f"- {point}\n"
            
            all_summaries.append({
                'file_path': os.path.abspath(file_path),
                'category': category,
                'summary': summary_text_for_file
            })

            if category in file_management_settings and file_management_settings[category].get('destination'):
                settings = file_management_settings[category]
                prefix = settings.get('prefix', '')
                destination_folder = settings['destination']
                
                os.makedirs(destination_folder, exist_ok=True)

                original_file_name, original_file_extension = os.path.splitext(os.path.basename(file_path))
                
                new_file_name = f"{prefix}{original_file_name}{original_file_extension}"
                destination_path = os.path.join(destination_folder, new_file_name)
                
                try:
                    shutil.move(file_path, destination_path)
                    print(f"Moved and renamed: {file_path} -> {destination_path}")
                except Exception as move_e:
                    print(f"Error moving original file: {move_e}")
            
        else:
            print("Step 8.1: Skipping file - unable to extract meaningful text.")
    
    print("\n--- Step 9: All supported files processed ---")
    return all_summaries

# --- Get Color Choice Function ---
def get_color_choice():
    """
    Prompts the user for a color choice and returns the corresponding ANSI code.
    """
    color_map = {
        'red': "\033[91m",
        'green': "\033[92m",
        'yellow': "\033[93m",
        'blue': "\033[94m",
        'magenta': "\033[95m",
        'cyan': "\033[96m"
    }

    print("\nChoose a color for bolded text:")
    for color in color_map:
        print(f"- {color}")
    
    choice = input("Enter your choice [default: yellow]: ").strip().lower()
    return color_map.get(choice, color_map['yellow'])


# --- Main Execution ---
if __name__ == "__main__":
    try:
        # Load last used models if they exist
        last_summarizer_model, last_classifier_model = load_last_models()

        # Step 0: Get the download directory for models and color choice
        cache_directory = get_download_directory()
        color_code = get_color_choice()

        # Question 1: Select the summarization model
        summarizer_model_name = select_summarization_model(last_summarizer_model)

        # Question 2: Ask if user wants to select a different classifier
        classifier_model_name = "MoritzLaurer/xtremedistil-l6-h256-mnli-fever-anli-ling-binary"
        classifier_choice = input(f"\nDo you want to select a different classifier model? (y/n) [default: n]: ").strip().lower()
        if classifier_choice in ['y', 'yes']:
            classifier_model_name = select_classifier_model(last_classifier_model)
        else:
            # If the user chose not to select a new model, and we have a last one, use it.
            if last_classifier_model:
                classifier_model_name = last_classifier_model
        
        # Question 3: Set the classifier threshold for the selected model
        classifier_threshold_to_use = get_classifier_threshold(classifier_model_name)
        
        # Question 4: Set summary lengths
        min_summary_length, max_summary_length = get_summary_lengths()
        
        # Question 5: Set the max token chunk size
        token_chunk_size = get_token_chunk_size()


        # --- Save last models for next run ---
        save_last_models(summarizer_model_name, classifier_model_name)


        # --- Load Models ---
        print(f"\nStep 1: Loading summarization model ({summarizer_model_name}) and classification model...")
        print("If this is the first time you are running the script, a large file download will begin now. Please wait for it to complete.")
        print("Note: If the script appears unresponsive during this step, it is likely downloading a large file. Forcing a stop with Ctrl+C may not be immediate during these operations.")

        try:
            print("Step 1a: Explicitly loading tokenizer...")
            summarizer_tokenizer = AutoTokenizer.from_pretrained(summarizer_model_name, cache_dir=cache_directory)
            
            print("Step 1b: Explicitly loading model...")
            summarizer_model = AutoModelForSeq2SeqLM.from_pretrained(summarizer_model_name, cache_dir=cache_directory).to(device)
            
            print("Step 1c: Creating summarization pipeline...")
            summarizer_pipeline = pipeline(
                "summarization",
                model=summarizer_model,
                tokenizer=summarizer_tokenizer,
                device=0 if device == 'cuda' else -1,
                framework="pt"
            )
            print("Step 1d: Summarization model loaded successfully.")
        except Exception as e:
            print(f"An error occurred while loading the summarization model: {e}")
            print("Please check the model name and your internet connection.")
            exit()

        try:
            classifier_pipeline = pipeline(
                task="zero-shot-classification",
                model=classifier_model_name,
                device=device,
                cache_dir=cache_directory
            )
            print(f"Step 1.1: Classification model ({os.path.basename(classifier_model_name)}) loaded successfully with custom threshold.")
        except Exception as e:
            print(f"Error loading classification model: {e}")
            classifier_pipeline = None
            print("Step 1.1: Falling back to 'Other' for all classifications.")

        CATEGORIES = load_and_edit_categories()

        while True:
            folder_path = input("Enter the folder path containing the files: ").strip()
            if os.path.isdir(folder_path):
                break
            print("\nError: The provided path is not a valid directory. Please try again.")

        scan_sub_input = input("Scan subdirectories? (y/n) [default: n]: ").strip().lower()
        scan_subdirectories = scan_sub_input in ['y', 'yes']

        while True:
            try:
                start_chunk_input = input("Enter the number of the first chunk (e.g., page number): ").strip()
                start_chunk_to_use = int(start_chunk_input)
                if start_chunk_to_use > 0:
                    break
                else:
                    print("Please enter a number greater than 0.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")
        
        while True:
            try:
                end_chunk_input = input("Enter the number of the last chunk: ").strip()
                end_chunk_to_use = int(end_chunk_input)
                if end_chunk_to_use >= start_chunk_to_use:
                    break
                else:
                    print("Please enter a number greater than or equal to the starting chunk number.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")

        file_management_settings = {}
        move_and_rename_choice = input("\nDo you want to move and rename files to subfolders based on their category? (y/n) [default: n]: ").strip().lower()
        
        if move_and_rename_choice in ['y', 'yes']:
            print("\n--- Define Global File Management Rules ---")
            prefix = input("Enter a global prefix for all renamed files (e.g., 'PROJ_') [leave blank for none]: ").strip()
            
            for category in CATEGORIES:
                if category != "Other":
                    dest_folder_name = category.replace(' ', '_')
                    file_management_settings[category] = {
                        'prefix': prefix,
                        'destination': os.path.join(folder_path, dest_folder_name)
                    }
            print("File management rules have been set up automatically for all categories.")


        print("\nStep 10: Starting the batch process...")
        
        all_summaries = process_files_in_folder(
            folder_path, 
            scan_subdirectories, 
            CATEGORIES, 
            start_chunk_to_use, 
            end_chunk_to_use, 
            file_management_settings,
            summarizer_pipeline,
            classifier_pipeline,
            min_summary_length,
            max_summary_length,
            classifier_threshold_to_use,
            token_chunk_size,
            color_code
        )

        if all_summaries:
            consolidated_summary_file = os.path.join(folder_path, "all_summaries.txt")
            print(f"\n--- Saving all summaries to {consolidated_summary_file} ---")
            with open(consolidated_summary_file, "w", encoding="utf-8") as f:
                for summary_data in all_summaries:
                    f.write("="*50 + "\n")
                    f.write(f"File: {summary_data['file_path']}\n")
                    f.write(f"Category: {summary_data['category']}\n\n")
                    f.write(f"Summary:\n{summary_data['summary']}\n")
            print("All summaries consolidated.")
            
    except KeyboardInterrupt:
        print("\nProcess interrupted by user. Exiting gracefully.")
    except Exception as e:
        print(f"\nAn unexpected error occurred: {e}")
