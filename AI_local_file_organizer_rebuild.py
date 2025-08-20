# comboaisort41.py
import os
import sys

# Suppress TensorFlow/oneDNN warnings
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '2'
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'

import fitz # PyMuPDF for PDF
from docx import Document
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook
from transformers import pipeline
import torch
from transformers import AutoTokenizer, AutoModelForSequenceClassification
import shutil
import json
import colorama
from colorama import Fore, Back, Style
import re
import multiprocessing

# Initialize colorama to auto-reset styles after each print
colorama.init(autoreset=True)

# Define a custom color for dark brown
DARK_BROWN = '\033[38;2;101;67;33m'

# Set global output line limit
MAX_OUTPUT_LINES = 5
MAX_LINE_WIDTH = 120
MAX_VERBIAGE_CHARS = 220
MAX_VERBIAGE_LINES = 5

# --- Global Variables for Models ---
summarizer = None
classifier = None

# --- Helper Function ---
def sanitize_for_path(name):
    """Truncates and sanitizes a string to be a valid directory/file name."""
    sanitized = re.sub(r'[<>:"/\\|?*]', '', name)
    return sanitized[:50].strip()

# --- Load and Edit Categories Function ---
def load_and_edit_categories():
    category_file = "categories.json"
    
    # New default categories based on user's prompt
    default_categories = [
        "a document related to cancer or medical oncology.",
        "a document related to modern global politics, international relations, or war, including foreign policy or political instability.",
        "a document related to public health, epidemiology, or population-level disease trends.",
        "a document related to the brain, spine, or neurology.",
        "a document related to the heart or cardiology.",
        "information about herbal remedies, vitamins, or dietary supplements.",
        "information regarding medical diagnostics, imaging, or surgical procedures.",
        "the principles or practices of traditional chinese medicine (tcm).",
        "a document about home improvement, gardening, or vehicle maintenance.",
        "a document about outdoor activities, sports, or recreation.",
        "a document about planetary science, astronomy, or solar system phenomena.",
        "a document related to geology, geography, or travel.",
        "a document about high-tech, software, or information technology, business, finance, or economics, including markets, investment, management, 

computer security, or money.",
        "a document about natural sciences or earth systems, including topics like geology, astronomy, or climate, planetary science, or solar system 

phenomena.",
        "a document related to archaeology, or human evolution, including topics like hominins, genetics, fossils, or ancient artifacts.",
        "a document related to ancient history, religion, esoteric beliefs, or mythology.",
        "a document about traditional or herbal medicine, including phyto-compounds or traditional chinese medicine (tcm).",
        "a document concerning genetics, hereditary diseases, or congenital conditions.",
        "a document related to cancer or medical oncology.",
        "a document related to human health and aging well, herbal remedies, vitamins, or dietary supplements.",
        "a document related to the brain, spine, or neurology.",
        "a document related to the circulatory system, including heart and blood, and arteries.",
        "a document about political ideologies, social dynamics, or historical conflicts, including radicalization, propaganda, or genocides, jewish 

history, the holocaust, or zionism, gaza, israeli-palestinian relations or middle eastern geopolitics.",
        "a document about home improvement, gardening, or vehicle maintenance.",
        "a document about outdoor activities, sports, or recreation.",
        "a document related to geography, or travel."
    ]
    categories = default_categories
    try:
        if os.path.exists(category_file):
            with open(category_file, 'r', encoding='utf-8') as f:
                categories = json.load(f)
            print(f"{Fore.CYAN}Step 1.2: Loaded categories from previous session.{Style.RESET_ALL}")
    except Exception as e:
        print(f"{Fore.LIGHTBLACK_EX}WARNING: Error loading categories file: {e}. Using default categories.{Style.RESET_ALL}")

    print(f"\n{Fore.LIGHTYELLOW_EX}--- Current Categories ---{Style.RESET_ALL}")
    for i, cat in enumerate(categories): print(f"{Fore.LIGHTYELLOW_EX}[{i+1}] {cat}{Style.RESET_ALL}")
    print(f"{Fore.LIGHTYELLOW_EX}--------------------------{Style.RESET_ALL}")

    prompt = "Do you want to edit these categories? (y/n) [default: no]: "
    user_input = input(f"{Fore.GREEN}{prompt}{Style.RESET_ALL}").strip().lower()
    if user_input in ['y', 'yes']:
        instructions = f"\n{Fore.LIGHTBLACK_EX}Editing categories: 'add <name>', 'remove <num>', 'edit <num> <name>', 'list', 'done'{Style.RESET_ALL}"
        print(instructions)
        while True:
            command = input("> ").strip().lower()
            parts = command.split()
            if not parts: continue
            action = parts[0]
            if action == 'add' and len(parts) >= 2: categories.append(" ".join(parts[1:])); print(f"{Fore.CYAN}Added: {' '.join(parts[1:])}

{Style.RESET_ALL}")
            elif action == 'remove' and len(parts) == 2 and parts[1].isdigit():
                idx = int(parts[1]) - 1
                if 0 <= idx < len(categories): print(f"{Fore.RED}Removed: {categories.pop(idx)}{Style.RESET_ALL}")
                else: print(f"{Fore.RED}Invalid number.{Style.RESET_ALL}")
            elif action == 'edit' and len(parts) >= 3 and parts[1].isdigit():
                idx = int(parts[1]) - 1
                if 0 <= idx < len(categories):
                    old_cat = categories[idx]
                    categories[idx] = " ".join(parts[2:])
                    print(f"{Fore.CYAN}Edited '{old_cat}' to '{categories[idx]}'{Style.RESET_ALL}")
                else: print(f"{Fore.RED}Invalid number.{Style.RESET_ALL}")
            elif action == 'list':
                for i, cat in enumerate(categories): print(f"{Fore.LIGHTWHITE_EX}[{i+1}] {cat}{Style.RESET_ALL}")
            else: print(f"{Fore.RED}Invalid command.{Style.RESET_ALL}")

    with open(category_file, 'w', encoding='utf-8') as f: json.dump(categories, f, indent=4)
    print(f"\n{Fore.CYAN}Step 1.3: Categories saved.{Style.RESET_ALL}")
    return categories

# --- Model and I/O Functions ---
def initialize_models():
    global summarizer, classifier
    print(f"\n{Fore.CYAN}Step 1: Loading AI models...{Style.RESET_ALL}")
    device = "cuda" if torch.cuda.is_available() else "cpu"
    print(f"{Fore.CYAN}Device set to use {'GPU' if device == 'cuda' else 'CPU'}.{Style.RESET_ALL}")
    summarizer = pipeline("summarization", model="t5-small", max_length=1024, device=device, framework="pt")
    try:
        classifier = pipeline("zero-shot-classification", model="MoritzLaurer/xtremedistil-l6-h256-zeroshot-v1.1-all-33", device=device, multi_label=True)
        print(f"{Fore.CYAN}Step 1.1: Classification model loaded.{Style.RESET_ALL}")
    except Exception as e:
        print(f"{Fore.LIGHTBLACK_EX}WARNING: Error loading classification model: {e}{Style.RESET_ALL}"); classifier = None

def run_summarization_in_process(summarizer_pipeline, chunk, queue, max_len, min_len):
    try:
        summary = summarizer_pipeline(chunk, max_length=max_len, min_length=min_len, truncation=True)
        queue.put(summary[0]['summary_text'])
    except Exception as e:
        queue.put(f"Summarization failed in subprocess: {e}")

def parse_chunks(chunks_string, file_path, max_chunk_length):
    """Parses a comma-separated string of chunks and ranges into a list of chunks."""
    if not chunks_string.strip():
        print(f"{Fore.LIGHTBLACK_EX}No chunks specified for {os.path.basename(file_path)}. Processing the first chunk.{Style.RESET_ALL}")
        return [1]

    chunks = set()
    parts = chunks_string.split(',')
    for part in parts:
        part = part.strip()
        if not part: continue
        if '-' in part:
            start_str, end_str = part.split('-')
            start, end = int(start_str), int(end_str)
            if start <= end:
                chunks.update(range(start, end + 1))
        else:
            chunks.add(int(part))
    return sorted(list(chunks))

def read_entire_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            try:
                with fitz.open(file_path) as doc:
                    for page in doc:
                        text += page.get_text() + "\n"
            except RuntimeError as e:
                print(f"{Fore.RED}ERROR: Failed to read PDF file '{os.path.basename(file_path)}' due to a format error: {e}. Skipping this file.

{Style.RESET_ALL}")
                return "__READ_ERROR__" # Use a special string to indicate a read error
        elif ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        elif ext == ".xlsx":
            for sheet in pd.ExcelFile(file_path).sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet)
                text += df.to_string(index=False) + "\n"
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text + "\n"
    except Exception as e:
        print(f"{Fore.LIGHTBLACK_EX}WARNING: Error reading {ext} file: {e}{Style.RESET_ALL}");
        return ""
    return text

def summarize_text(text, max_chunk_length, max_summary_length, min_summary_length, timeout_seconds, show_verbiage):
    if not text.strip():
        return ["No text to summarize."]

    words = text.split()
    chunks = [" ".join(words[i:i + max_chunk_length]) for i in range(0, len(words), max_chunk_length)]
    summaries = []

    for i, chunk in enumerate(chunks):
        q = multiprocessing.Queue()
        p = multiprocessing.Process(target=run_summarization_in_process, args=(summarizer, chunk, q, max_summary_length, min_summary_length))
        p.start()
        p.join(timeout=timeout_seconds)

        if p.is_alive():
            print(f"{Fore.RED}Step 7.1.3: Summarization for chunk {i + 1} timed out after {timeout_seconds} seconds.{Style.RESET_ALL}")
            p.terminate()
            p.join()
            summaries.append("Summary failed due to timeout.")
        else:
            try:
                result = q.get(timeout=5)
                summaries.append(result)
                print(f"{Fore.CYAN}Step 7.1.2: Summarization for chunk {i + 1} complete.{Style.RESET_ALL}")
            except Exception:
                summaries.append("Summary failed to return result from process.")
                print(f"{Fore.LIGHTBLACK_EX}WARNING: Error retrieving result for chunk {i + 1}.{Style.RESET_ALL}")

    return ". ".join(summaries).split(". ")

def categorize_summary(summary_text, categories):
    if classifier is None: return {'labels': ['Other'], 'scores': [1.0]}
    print(f"{Fore.CYAN}Step 8.1: Sending summary to classifier model...{Style.RESET_ALL}")
    try:
        results = classifier(summary_text, candidate_labels=categories, multi_label=True)
        print(f"{Fore.GREEN}Step 8.2: Categorization complete.{Style.RESET_ALL}")
        return results
    except Exception as e:
        print(f"{Fore.LIGHTBLACK_EX}WARNING: Error during AI categorization: {e}{Style.RESET_ALL}"); return {'labels': ['Other'], 'scores': [1.0]}

def tag_file_properties(file_path, primary_category, secondary_categories):
    if not os.access(file_path, os.W_OK):
        print(f"{Fore.RED}Tagging skipped: No write permission for '{os.path.basename(file_path)}'.{Style.RESET_ALL}")
        return
    ext = os.path.splitext(file_path)[1].lower()
    if not os.path.exists(file_path):
        print(f"{Fore.LIGHTBLACK_EX}WARNING: Tagging skipped: File not found at {file_path}{Style.RESET_ALL}"); return
    keyword_parts = []
    if primary_category and primary_category != "Other": keyword_parts.append(f"AILFO Primary:; {primary_category}")
    for cat in secondary_categories: keyword_parts.append(f"2nd:; {cat}")
    keyword_string = ":: ".join(keyword_parts)
    if not keyword_string:
        print(f"{Fore.LIGHTBLACK_EX}No valid categories to tag for {os.path.basename(file_path)}{Style.RESET_ALL}"); return
    try:
        obj = None
        if ext == ".pdf":
            doc = fitz.open(file_path)
            doc.metadata["keywords"] = keyword_string
            doc.save(file_path, incremental=True, encryption=fitz.PDF_ENCRYPT_KEEP); doc.close()
        elif ext in [".docx", ".pptx", ".xlsx"]:
            if ext == ".docx": obj = Document(file_path)
            elif ext == ".pptx": obj = Presentation(file_path)
            elif ext == ".xlsx": obj = load_workbook(file_path)
            if ext != ".xlsx": obj.core_properties.keywords = keyword_string
            else: obj.properties.keywords = keyword_string
            obj.save(file_path)
        else:
            print(f"{Fore.LIGHTBLACK_EX}Tagging not supported for {ext} files.{Style.RESET_ALL}"); return
        
        colored_output = f"{Back.BLACK}{Fore.WHITE}Tagged '{os.path.basename(file_path)}' with: {keyword_string.replace('2nd:;', f'{Fore.YELLOW}2nd:;

{Fore.WHITE}')}{Style.RESET_ALL}"
        print(colored_output)
    except Exception as e:
        print(f"{Fore.LIGHTBLACK_EX}WARNING: Error tagging file {os.path.basename(file_path)}: {e}{Style.RESET_ALL}")

def create_shortcut(source_path, shortcut_path):
    if sys.platform == "win32":
        shortcut_path += ".url"
        if os.path.exists(shortcut_path): return
        with open(shortcut_path, "w") as f:
            f.write(f"[InternetShortcut]\nURL=file:///{os.path.abspath(source_path)}\n")
    else:
        if os.path.lexists(shortcut_path): return
        os.symlink(source_path, shortcut_path)

def preprocess_by_filename(file_list, categories, file_action, output_folder_path, file_management_settings, settings):
    print(f"\n{Fore.CYAN}--- Starting Pre-processing Step (Filename and Content Scan) ---{Style.RESET_ALL}")
    remaining_files = list(file_list)
    for file_path in file_list:
        file_name_no_ext = os.path.splitext(os.path.basename(file_path))[0]
        words_in_name = {word.lower() for word in file_name_no_ext.replace('_', ' ').replace('-', ' ').split()}
        found_category = None
        for cat in categories:
            if set(cat.lower().split()).issubset(words_in_name):
                found_category = cat; break
        if found_category:
            sanitized_cat = sanitize_for_path(found_category)
            destination_folder = os.path.join(output_folder_path, sanitized_cat)
            os.makedirs(destination_folder, exist_ok=True)
            prefix = file_management_settings.get(found_category, {}).get('prefix', '')
            new_file_name = f"{prefix}_{file_name_no_ext}" if prefix else file_name_no_ext
            
            destination_path = os.path.join(destination_folder, new_file_name)
            
            if file_action != 'shortcut':
                destination_path_with_ext = destination_path + os.path.splitext(file_path)[1]
            else:
                destination_path_with_ext = destination_path
            
            try:
                if file_action == 'shortcut':
                    create_shortcut(file_path, destination_path)
                    print(f"{Fore.CYAN}Pre-processed shortcut for: {os.path.basename(file_path)} -> {destination_folder}{Style.RESET_ALL}")
                else:
                    if file_action == 'move':
                        shutil.move(file_path, destination_path_with_ext)
                        print(f"{Fore.CYAN}Pre-processed and moved: {os.path.basename(file_path)} -> {destination_folder}{Style.RESET_ALL}")
                    elif file_action == 'copy':
                        shutil.copy2(file_path, destination_path_with_ext)
                        print(f"{Fore.CYAN}Pre-processed and copied: {os.path.basename(file_path)} -> {destination_folder}{Style.RESET_ALL}")
                remaining_files.remove(file_path)
            except Exception as e:
                print(f"{Fore.LIGHTBLACK_EX}WARNING: Error pre-processing {os.path.basename(file_path)}: {e}{Style.RESET_ALL}")
    print(f"{Fore.CYAN}--- Pre-processing Complete: {len(file_list) - len(remaining_files)} files processed. ---{Style.RESET_ALL}")
    return remaining_files

def extract_keywords_from_categories(categories):
    keywords = set()
    for category in categories:
        words = re.split(r'[,;.\s]\s*', category.lower())
        stop_words = {'a', 'an', 'and', 'about', 'with', 'in', 'or', 'to', 'the', 'of', 'related', 'document', 'including', 'information', 'for', 'like'}
        keywords.update(word for word in words if word and word not in stop_words)
    return list(keywords)

def filter_text_by_keywords(text, keywords, before_chars, after_chars):
    if not keywords:
        return text
    
    filtered_snippets = []
    
    pattern = re.compile(r'\b(' + '|'.join(re.escape(k) for k in keywords) + r')\b', re.IGNORECASE)

    for match in pattern.finditer(text):
        start_index = max(0, match.start() - before_chars)
        end_index = min(len(text), match.end() + after_chars)
        
        start_bound = text.rfind(' ', 0, start_index)
        end_bound = text.find(' ', end_index)
        
        start_index = start_bound + 1 if start_bound != -1 else start_index
        end_index = end_bound if end_bound != -1 else end_index

        snippet = text[start_index:end_index].strip()
        if snippet:
            filtered_snippets.append(snippet)
            
    return "... ".join(filtered_snippets)

def cull_to_complete_sentences(text):
    sentences = re.split(r'(?<=[.!?])\s+', text)
    culled_sentences = [s for s in sentences if s.strip().endswith(('.', '!', '?'))]
    return " ".join(culled_sentences)

def process_single_file(task_data):
    file_path, file_path_index, total_files, categories, file_management_settings, settings = task_data
    
    print(f"\n{Back.BLACK}{Fore.RED}File: {file_path}{Style.RESET_ALL}")
    
    file_name_no_ext = os.path.splitext(os.path.basename(file_path))[0]
    cleaned_file_name = file_name_no_ext.replace('_', ' ').replace('-', ' ')
    
    should_skip = False
    
    text = read_entire_file(file_path)
    
    if text == "__READ_ERROR__":
        special_category = "unreadable file"
        sanitized_cat = sanitize_for_path(special_category)
        destination_folder = os.path.join(settings['output_folder_path'], sanitized_cat)
        os.makedirs(destination_folder, exist_ok=True)
        
        if settings['file_action_choice'] == 'shortcut':
            destination_path_with_ext = os.path.join(destination_folder, cleaned_file_name) + ".url"
        else:
            destination_path_with_ext = os.path.join(destination_folder, cleaned_file_name + os.path.splitext(file_path)[1])
        
        if os.path.exists(destination_path_with_ext):
            print(f"{Fore.LIGHTBLACK_EX}Skipping: Output file '{os.path.basename(destination_path_with_ext)}' already exists in '{sanitized_cat}' 

subfolder.{Style.RESET_ALL}")
            return None
            
        print(f"{Fore.YELLOW}WARNING: File format error detected. Moving to '{special_category}' folder.{Style.RESET_ALL}")
        
        try:
            if settings['file_action_choice'] == 'shortcut':
                create_shortcut(file_path, destination_path_with_ext)
                print(f"{Fore.CYAN}Created shortcut in -> {destination_folder}{Style.RESET_ALL}")
            elif settings['file_action_choice'] == 'move':
                shutil.move(file_path, destination_path_with_ext)
                print(f"{Fore.CYAN}Moved and renamed: {file_path} -> {destination_path_with_ext}{Style.RESET_ALL}")
            elif settings['file_action_choice'] == 'copy':
                shutil.copy2(file_path, destination_path_with_ext)
                print(f"{Fore.CYAN}Copied and renamed: {file_path} -> {destination_path_with_ext}{Style.RESET_ALL}")
        except Exception as e:
            print(f"{Fore.LIGHTBLACK_EX}Error during file operation for '{special_category}': {e}{Style.RESET_ALL}")

        return {'file_path': os.path.abspath(file_path), 'primary_category': special_category, 'all_categories': [special_category], 'summary': 'File 

could not be read due to format error.'}
    
    if not text.strip():
        special_category = "contains no text"
        sanitized_cat = sanitize_for_path(special_category)
        destination_folder = os.path.join(settings['output_folder_path'], sanitized_cat)
        
        if settings['file_action_choice'] == 'shortcut':
            destination_path_with_ext = os.path.join(destination_folder, cleaned_file_name) + ".url"
        else:
            destination_path_with_ext = os.path.join(destination_folder, cleaned_file_name + os.path.splitext(file_path)[1])
        
        if os.path.exists(destination_path_with_ext):
            print(f"{Fore.LIGHTBLACK_EX}Skipping: Output file '{os.path.basename(destination_path_with_ext)}' already exists in '{sanitized_cat}' 

subfolder.{Style.RESET_ALL}")
            return None
            
        os.makedirs(destination_folder, exist_ok=True)
        print(f"{Fore.YELLOW}WARNING: No text found in file. Moving to 'contains no text' folder.{Style.RESET_ALL}")
        
        try:
            if settings['file_action_choice'] == 'shortcut':
                create_shortcut(file_path, os.path.join(destination_folder, cleaned_file_name))
                print(f"{Fore.CYAN}Created shortcut in -> {destination_folder}{Style.RESET_ALL}")
            elif settings['file_action_choice'] == 'move':
                shutil.move(file_path, destination_path_with_ext)
                print(f"{Fore.CYAN}Moved and renamed: {file_path} -> {destination_path_with_ext}{Style.RESET_ALL}")
            elif settings['file_action_choice'] == 'copy':
                shutil.copy2(file_path, destination_path_with_ext)
                print(f"{Fore.CYAN}Copied and renamed: {file_path} -> {destination_path_with_ext}{Style.RESET_ALL}")
        except Exception as e:
            print(f"{Fore.LIGHTBLACK_EX}Error during file operation for '{special_category}': {e}{Style.RESET_ALL}")

        return {'file_path': os.path.abspath(file_path), 'primary_category': special_category, 'all_categories': [special_category], 'summary': 'No text 

found in file.'}
    
    should_exclude = False
    negative_keywords_str = settings.get('negative_filter_keywords', '')
    negative_keywords = [k.strip().lower() for k in negative_keywords_str.split(',') if k.strip()]
    
    if negative_keywords:
        for keyword in negative_keywords:
            if keyword in text.lower():
                print(f"{Fore.RED}Skipping '{os.path.basename(file_path)}' due to negative keyword: '{keyword}'{Style.RESET_ALL}")
                should_exclude = True
                break
    
    exclude_regex_str = settings.get('exclude_regex', '')
    if not should_exclude and exclude_regex_str:
        try:
            exclude_regex = re.compile(exclude_regex_str, re.IGNORECASE)
            if exclude_regex.search(text):
                print(f"{Fore.RED}Skipping '{os.path.basename(file_path)}' due to exclusion regex match.{Style.RESET_ALL}")
                should_exclude = True
        except re.error as e:
            print(f"{Fore.RED}ERROR: Invalid regex string for exclusion: {e}. Skipping regex check for this session.{Style.RESET_ALL}")
    
    if should_exclude:
        return None

    if settings['file_action_choice'] != 'none':
        already_processed = False
        for category in categories:
            if category == "Other": continue
            prefix = file_management_settings.get(category, {}).get('prefix', '')
            new_file_name_no_ext = f"{prefix}_{cleaned_file_name}" if prefix else cleaned_file_name
            if settings['file_action_choice'] == 'shortcut':
                expected_output_name = new_file_name_no_ext + ".url" if sys.platform == "win32" else new_file_name_no_ext
            else:
                expected_output_name = new_file_name_no_ext + os.path.splitext(file_path)[1]

            if expected_output_name in existing_files_cache:
                print(f"{Fore.LIGHTBLACK_EX}Skipping pre-filtered: '{os.path.basename(file_path)}' found in '{sanitize_for_path(category)}' subfolder.

{Style.RESET_ALL}")
                already_processed = True
                break
                
        for special_cat in ["contains no text", "unreadable file"]:
            sanitized_cat = sanitize_for_path(special_cat)
            if settings['file_action_choice'] == 'shortcut':
                expected_output_name = cleaned_file_name + ".url" if sys.platform == "win32" else cleaned_file_name
            else:
                expected_output_name = cleaned_file_name + os.path.splitext(file_path)[1]

            if expected_output_name in existing_files_cache:
                print(f"{Fore.LIGHTBLACK_EX}Skipping pre-filtered: '{os.path.basename(file_path)}' found in '{sanitized_cat}' subfolder.

{Style.RESET_ALL}")
                should_skip = True
                break

        if not should_skip:
            filtered_files.append(file_path)
    
    files_to_process = filtered_files
    
    print(f"{Fore.CYAN}Pre-filtering complete: {len(files_to_process)} new files to process out of a total of {total_initial_files} from the initial 

list.{Style.RESET_ALL}")

    is_pdf = os.path.splitext(file_path)[1].lower() == '.pdf'
    is_text = os.path.splitext(file_path)[1].lower() == '.txt'
    
    max_chunk_length = settings['generic_chunk_length']
    chunks_list = []
    
    if is_pdf:
        try:
            with fitz.open(file_path) as doc:
                page_count = doc.page_count
                if page_count > settings['long_pdf_threshold']:
                    max_chunk_length = settings['long_pdf_chunk_length']
                    chunks_list = parse_chunks(settings['long_pdf_chunks_list'], file_path, max_chunk_length)
                    print(f"{Fore.CYAN}Detected long PDF ({page_count} pages). Using long PDF settings.{Style.RESET_ALL}")
                else:
                    max_chunk_length = settings['short_pdf_chunk_length']
                    chunks_list = parse_chunks(settings['short_pdf_chunks_list'], file_path, max_chunk_length)
                    print(f"{Fore.CYAN}Detected short PDF ({page_count} pages). Using short PDF settings.{Style.RESET_ALL}")
        except Exception as e:
            print(f"{Fore.LIGHTBLACK_EX}WARNING: Error getting PDF page count: {e}. Using default PDF settings.{Style.RESET_ALL}")
            max_chunk_length = settings['pdf_chunk_length']
            chunks_list = parse_chunks(settings['pdf_chunks_list'], file_path, max_chunk_length)
    elif is_text:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
                word_count = len(content.split())
                if word_count > settings['long_text_threshold']:
                    max_chunk_length = settings['long_text_chunk_length']
                    chunks_list = parse_chunks(settings['long_text_chunks_list'], file_path, max_chunk_length)
                    print(f"{Fore.CYAN}Detected long text file ({word_count} words). Using long text settings.{Style.RESET_ALL}")
                else:
                    max_chunk_length = settings['short_text_chunk_length']
                    chunks_list = parse_chunks(settings['short_text_chunks_list'], file_path, max_chunk_length)
                    print(f"{Fore.CYAN}Detected short text file ({word_count} words). Using short text settings.{Style.RESET_ALL}")
        except Exception as e:
            print(f"{Fore.LIGHTBLACK_EX}WARNING: Error counting words in text file: {e}. Using default text settings.{Style.RESET_ALL}")
            max_chunk_length = settings['generic_chunk_length']
            chunks_list = parse_chunks(settings['generic_chunks_list'], file_path, max_chunk_length)
    else:
        max_chunk_length = settings['generic_chunk_length']
        chunks_list = parse_chunks(settings['generic_chunks_list'], file_path, max_chunk_length)

    if not chunks_list:
        print(f"{Fore.LIGHTBLACK_EX}WARNING: No chunks specified for this file type and size. Skipping.{Style.RESET_ALL}")
        return None
    
    chunk_string = ", ".join(map(str, chunks_list))
    print(f"{Fore.CYAN}Step 6: Extracting from chunk(s): {chunk_string}...{Style.RESET_ALL}")
    text = read_entire_file(file_path)
    if not text.strip():
        print(f"{Fore.LIGHTBLACK_EX}WARNING: No text found in file, skipping.{Style.RESET_ALL}")
        return None
    print(f"{Fore.CYAN}Step 6.1: Text extracted successfully. {len(text.split())} words found.{Style.RESET_ALL}")
    
    print(f"{Fore.CYAN}Step 7: Pre-processing text for summarization...{Style.RESET_ALL}")

    text_to_summarize = ""
    culling_mode = settings.get('culling_mode', 'c')
    
    max_culling_words = settings['generic_chunk_length']
    if is_pdf:
        try:
            with fitz.open(file_path) as doc:
                page_count = doc.page_count
                if page_count > settings['long_pdf_threshold']:
                    max_culling_words = settings['long_pdf_chunk_length']
                else:
                    max_culling_words = settings['short_pdf_chunk_length']
        except Exception:
            pass
    elif is_text:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                word_count = len(f.read().split())
                if word_count > settings['long_text_threshold']:
                    max_culling_words = settings['long_text_chunk_length']
                else:
                    max_culling_words = settings['short_text_chunk_length']
        except Exception:
            pass
            
    has_keywords = False
    filter_keywords = [k.strip().lower() for k in settings.get('filter_keywords', '').split(',') if k.strip()]
    if filter_keywords:
        for keyword in filter_keywords:
            if re.search(r'\b' + re.escape(keyword) + r'\b', text, re.IGNORECASE):
                has_keywords = True
                break

    print(f"{Fore.CYAN}Culling Mode: {culling_mode.upper()}. ", end="")

    if culling_mode == 'r':
        if settings.get('culling_regex'):
            print(f"Using provided regex for culling.{Style.RESET_ALL}")
            try:
                culling_regex = re.compile(settings['culling_regex'], re.DOTALL)
                matches = culling_regex.findall(text)
                text_to_summarize = " ".join(matches)
            except re.error as e:
                print(f"{Fore.RED}ERROR: Invalid regex string provided for culling: {e}. Falling back to default culling.{Style.RESET_ALL}")
                text_to_summarize = cull_to_complete_sentences(text)
        else:
            print(f"No culling regex provided. Culling text to complete sentences.{Style.RESET_ALL}")
            text_to_summarize = cull_to_complete_sentences(text)
    elif culling_mode in ['s', 'b', '2', 'p']:
        if not filter_keywords:
            print(f"No keywords provided. Using default sentence culling.{Style.RESET_ALL}")
            text_to_summarize = cull_to_complete_sentences(text)
        else:
            if has_keywords:
                print(f"{Fore.GREEN}Keyword(s) found. Culling successful.{Style.RESET_ALL}")
                if culling_mode == 's':
                    sentence_pattern = r'[^.!?]*\b(' + '|'.join(re.escape(k) for k in filter_keywords) + r')\b[^.!?]*[.!?]'
                    matches = re.findall(sentence_pattern, text, re.IGNORECASE | re.DOTALL)
                    text_to_summarize = " ".join(matches)
                elif culling_mode == 'b':
                    pattern = r'((?:[^.!?]*[.!?]\s+){0,1}[^.!?]*\b(' + '|'.join(re.escape(k) for k in filter_keywords) + r')\b[^.!?]*[.!?])'
                    matches = re.findall(pattern, text, re.IGNORECASE | re.DOTALL)
                    text_to_summarize = " ".join([m[0] for m in matches])
                elif culling_mode == '2':
                    pattern = r'((?:[^.!?]*[.!?]\s+){0,2}[^.!?]*\b(' + '|'.join(re.escape(k) for k in filter_keywords) + r')\b[^.!?]*[.!?](?:\s+[^.!?]*

[.!?]){0,2})'
                    matches = re.findall(pattern, text, re.IGNORECASE | re.DOTALL)
                    text_to_summarize = " ".join([m[0] for m in matches])
                elif culling_mode == 'p':
                    pattern = r'((?:[^\n]+\n){0,1}[^.!?]*\b(' + '|'.join(re.escape(k) for k in filter_keywords) + r')\b[^.!?]*(?:\s+[^\n]+\n){0,1})'
                    matches = re.findall(pattern, text, re.IGNORECASE | re.DOTALL)
                    text_to_summarize = " ".join([m[0] for m in matches])
            
                if not text_to_summarize.strip():
                    print(f"{Fore.YELLOW}WARNING: Culling produced no text. Falling back to summarizing entire extracted text.{Style.RESET_ALL}")
                    text_to_summarize = text
            else:
                print(f"{Fore.YELLOW}WARNING: No keywords found. Falling back to using defined chunks.{Style.RESET_ALL}")
                words = text.split()
                text_chunks = [" ".join(words[i:i + max_chunk_length]) for i in range(0, len(words), max_chunk_length)]
                
                chunks_to_use = []
                for chunk_num in chunks_list:
                    if 0 < chunk_num <= len(text_chunks):
                        chunks_to_use.append(text_chunks[chunk_num-1])
                text_to_summarize = " ".join(chunks_to_use)
                
    else: # culling_mode == 'n' or invalid choice
        print(f"No culling selected. Summarizing entire extracted text.{Style.RESET_ALL}")
        text_to_summarize = text

    if text_to_summarize.strip():
        words_culled = text_to_summarize.split()
        if len(words_culled) > max_culling_words:
            first_keyword_index = -1
            filter_keywords = [k.strip().lower() for k in settings.get('filter_keywords', '').split(',') if k.strip()]
            for word_index, word in enumerate(words_culled):
                if any(kw in word.lower() for kw in filter_keywords):
                    first_keyword_index = word_index
                    break
            
            if first_keyword_index != -1:
                start_index = max(0, first_keyword_index - max_culling_words // 2)
                end_index = min(len(words_culled), start_index + max_culling_words)
                trimmed_words = words_culled[start_index:end_index]
            else:
                trimmed_words = words_culled[:max_culling_words]
            
            text_to_summarize = " ".join(trimmed_words) + "..."
            
    if settings['show_chunk_verbiage'].lower() == 'y':
        raw_lines = text.split('\n')
        processed_lines = []
        temp_line = ""
        for line in raw_lines:
            words_in_line = line.strip().split()
            if len(words_in_line) < 3 and temp_line:
                temp_line += " | " + " ".join(words_in_line)
            elif len(words_in_line) < 3:
                temp_line = " ".join(words_in_line)
            else:
                if temp_line:
                    processed_lines.append(temp_line)
                temp_line = ""
                processed_lines.append(line.strip())
            
            if len(processed_lines) >= MAX_OUTPUT_LINES:
                break
        
        if temp_line and len(processed_lines) < MAX_OUTPUT_LINES:
            processed_lines.append(temp_line)

        display_text = '\n'.join(processed_lines[:MAX_OUTPUT_LINES])
        
        print(f"\n{Fore.MAGENTA}--- Raw Text Verbiage ({len(text.split())} words) ---{Style.RESET_ALL}")
        lines_to_display = display_text.split('\n')
        
        for line_idx, line in enumerate(lines_to_display):
            if line_idx >= MAX_VERBIAGE_LINES:
                print(f"{Fore.MAGENTA}...{Style.RESET_ALL}")
                break
            
            wrapped_line = '\n'.join([line[j:j+MAX_LINE_WIDTH] for j in range(0, len(line), MAX_LINE_WIDTH)])
            
            if len(wrapped_line) > MAX_VERBIAGE_CHARS:
                wrapped_line = wrapped_line[:MAX_VERBIAGE_CHARS] + "..."
            
            filter_keywords = [k.strip().lower() for k in settings.get('filter_keywords', '').split(',') if k.strip()]
            for keyword in filter_keywords:
                wrapped_line = re.sub(r'(?i)\b' + re.escape(keyword) + r'\b', f"{Fore.YELLOW}{Style.BRIGHT}\\g<0>{Fore.MAGENTA}{Style.NORMAL}", 

wrapped_line)
            print(f"{Fore.MAGENTA}{wrapped_line}{Style.RESET_ALL}")
            
        print(f"{Fore.MAGENTA}------------------------------{Style.RESET_ALL}\n")

    if not text_to_summarize.strip():
        print(f"{Fore.LIGHTBLACK_EX}WARNING: No text for summarization, skipping.{Style.RESET_ALL}")
        return None
    
    print(f"{Fore.CYAN}Step 7.1: Filtered text ready. Starting summarization for {Fore.LIGHTRED_EX}{os.path.basename(file_path)}{Style.RESET_ALL}

{Fore.CYAN}...{Style.RESET_ALL}")
    
    if settings['show_chunk_verbiage'].lower() == 'y':
        raw_lines = text_to_summarize.split('\n')
        processed_lines = []
        temp_line = ""
        for line in raw_lines:
            words_in_line = line.strip().split()
            if len(words_in_line) < 3 and temp_line:
                temp_line += " | " + " ".join(words_in_line)
            elif len(words_in_line) < 3:
                temp_line = " ".join(words_in_line)
            else:
                if temp_line:
                    processed_lines.append(temp_line)
                temp_line = ""
                processed_lines.append(line.strip())
            
            if len(processed_lines) >= MAX_OUTPUT_LINES:
                break
        
        if temp_line and len(processed_lines) < MAX_OUTPUT_LINES:
            processed_lines.append(temp_line)

        display_text = '\n'.join(processed_lines[:MAX_OUTPUT_LINES])
        
        print(f"\n{Fore.MAGENTA}--- Chunk 1 Verbiage ({len(text_to_summarize.split())} words) ---{Style.RESET_ALL}")
        lines_to_display = display_text.split('\n')
        
        for line_idx, line in enumerate(lines_to_display):
            if line_idx >= MAX_VERBIAGE_LINES:
                print(f"{Fore.MAGENTA}...{Style.RESET_ALL}")
                break
            
            wrapped_line = '\n'.join([line[j:j+MAX_LINE_WIDTH] for j in range(0, len(line), MAX_LINE_WIDTH)])
            
            if len(wrapped_line) > MAX_VERBIAGE_CHARS:
                wrapped_line = wrapped_line[:MAX_VERBIAGE_CHARS] + "..."
            
            filter_keywords = [k.strip().lower() for k in settings.get('filter_keywords', '').split(',') if k.strip()]
            for keyword in filter_keywords:
                wrapped_line = re.sub(r'(?i)\b' + re.escape(keyword) + r'\b', f"{Fore.YELLOW}{Style.BRIGHT}\\g<0>{Fore.MAGENTA}{Style.NORMAL}", 

wrapped_line)
            print(f"{Fore.MAGENTA}{wrapped_line}{Style.RESET_ALL}")
        
        print(f"{Fore.MAGENTA}------------------------------{Style.RESET_ALL}\n")
    
    print(f"{Fore.CYAN}Step 7.1.1: Summarizing chunk 1 for {Fore.LIGHTRED_EX}{os.path.basename(file_path)}{Style.RESET_ALL}{Fore.CYAN}...

{Style.RESET_ALL}")
    
    long_line = " ".join(text_to_summarize.split('\n'))
    lines_to_print = []
    start_idx = 0
    while start_idx < len(long_line) and len(lines_to_print) < MAX_OUTPUT_LINES:
        end_idx = start_idx + MAX_LINE_WIDTH
        if end_idx < len(long_line):
            last_space = long_line.rfind(' ', start_idx, end_idx)
            if last_space > start_idx:
                end_idx = last_space
        
        lines_to_print.append(long_line[start_idx:end_idx].strip())
        start_idx = end_idx + 1
        
    for line_idx, line in enumerate(lines_to_print):
        print(f"{Fore.CYAN}{line}{Style.RESET_ALL}")
    
    if len(lines_to_print) >= MAX_OUTPUT_LINES and len(long_line) > start_idx:
        print(f"{Fore.CYAN}...{Style.RESET_ALL}")

    max_summary_length_final = settings['max_summary_length']
    min_summary_length_final = settings['min_summary_length']
    if is_pdf:
        try:
            with fitz.open(file_path) as doc:
                page_count = doc.page_count
                if page_count > settings['long_pdf_threshold']:
                    max_summary_length_final = settings['long_pdf_max_summary_length']
                    min_summary_length_final = settings['long_pdf_min_summary_length']
                else:
                    max_summary_length_final = settings['short_pdf_max_summary_length']
                    min_summary_length_final = settings['short_pdf_min_summary_length']
        except Exception as e:
            print(f"{Fore.LIGHTBLACK_EX}WARNING: Error getting PDF page count: {e}. Using default summary settings.{Style.RESET_ALL}")
    
    bullet_points = summarize_text(text_to_summarize, max_chunk_length, max_summary_length_final, min_summary_length_final, settings

['summarizer_timeout'], settings['show_chunk_verbiage'])
    if not bullet_points or "Summary failed" in bullet_points[0]: print(f"{Fore.LIGHTBLACK_EX}WARNING: Summarization failed.{Style.RESET_ALL}"); return 

None

    print(f"\n{Fore.GREEN}--- Plain Text Summary ---{Style.RESET_ALL}")
    summary_text_for_file = ""
    for i, point in enumerate(bullet_points):
        if i >= MAX_OUTPUT_LINES:
            print(f"{Fore.LIGHTGREEN_EX}...{Style.RESET_ALL}")
            break
        p_strip = point.strip()
        if p_strip:
            wrapped_line = '\n'.join([p_strip[j:j+MAX_LINE_WIDTH] for j in range(0, len(p_strip), MAX_LINE_WIDTH)])
            print(f"{Fore.LIGHTGREEN_EX}- {wrapped_line}{Style.RESET_ALL}")
            summary_text_for_file += f"- {p_strip}\n"

    cleaned_file_name = os.path.splitext(os.path.basename(file_path))[0].replace('_', ' ').replace('-', ' ')
    
    results = categorize_summary(". ".join([cleaned_file_name] + bullet_points), categories)
    original_primary_category = results['labels'][0]
    original_primary_score = results['scores'][0]
    
    primary_category, secondary_categories, all_valid_categories = "Other", [], []
    if original_primary_score * 100 >= settings['confidence_threshold']:
        primary_category = original_primary_category
        all_valid_categories.append(primary_category)
        
        if len(results['labels']) > 1 and results['scores'][1] * 100 >= settings['secondary_confidence_threshold']:
            secondary_categories.append((results['labels'][1], results['scores'][1]))
            all_valid_categories.append(results['labels'][1])

    print(f"{Fore.CYAN}Step 9: Final summary complete.{Style.RESET_ALL}")

    if settings['tagging_enabled_choice'] in ['y', 'yes'] and all_valid_categories:
        sec_cat_names = [cat for cat, score in secondary_categories]
        tag_file_properties(file_path, primary_category, sec_cat_names)

    print(f"\n{Fore.GREEN}File: {Fore.LIGHTWHITE_EX}{file_path}{Style.RESET_ALL}")
    
    category_color = Fore.YELLOW if primary_category == "Other" else Fore.GREEN
    if primary_category == "Other" and original_primary_category != "Other":
        print(f"{Fore.GREEN}Primary Category:{Style.RESET_ALL} {category_color}{Style.BRIGHT}{primary_category} ({Fore.LIGHTBLACK_EX}was 

'{original_primary_category}' at {original_primary_score*100:.2f}%{Style.RESET_ALL}){Style.RESET_ALL}")
    else:
        print(f"{Fore.GREEN}Primary Category:{Style.RESET_ALL} {category_color}{Style.BRIGHT}{primary_category} ({Fore.LIGHTWHITE_EX}{results['scores']

[0]*100:.2f}%{Style.RESET_ALL})")

    if secondary_categories:
        colors = [Fore.CYAN, Fore.LIGHTGREEN_EX]
        colored_cats = [f"{colors[i % 2]}{cat} ({Fore.YELLOW}{score*100:.2f}%)" for i, (cat, score) in enumerate(secondary_categories)]
        print(f"{Fore.GREEN}Secondary Categories:{Style.RESET_ALL} {', '.join(colored_cats)}")

    print(f"\n{Fore.CYAN}--- Colorized Summary ---{Style.RESET_ALL}")
    if primary_category == "Other":
        display_summary = '\n'.join(summary_text_for_file.split('\n')[:MAX_OUTPUT_LINES])
        for line in display_summary.split('\n'):
            wrapped_line = '\n'.join([line[j:j+MAX_LINE_WIDTH] for j in range(0, len(line), MAX_LINE_WIDTH)])
            print(f"{Fore.MAGENTA}{wrapped_line}{Style.RESET_ALL}")
        return {'file_path': os.path.abspath(file_path), 'primary_category': primary_category, 'all_categories': all_valid_categories, 'summary': 

summary_text_for_file}

    sec_cat_map = {cat.lower(): i for i, (cat, score) in enumerate(secondary_categories)}
    colors = [Fore.CYAN, Fore.LIGHTGREEN_EX]
    
    display_lines = []
    for i, point in enumerate(bullet_points):
        if i >= MAX_OUTPUT_LINES:
            print(f"{Fore.LIGHTCYAN_EX}...{Style.RESET_ALL}")
            break
        p_strip = point.strip()
        if not p_strip: continue
        colorized_point = ""
        words_and_delimiters = re.split(r'([ ,.;:!?])', p_strip)
        for word in words_and_delimiters:
            if word.lower() == primary_category.lower(): colorized_point += f"{Fore.GREEN}{Style.BRIGHT}{word}{Style.RESET_ALL}"
            elif word.lower() in sec_cat_map:
                idx = sec_cat_map[word.lower()]
                color = colors[idx % 2]
                colorized_point += f"{color}{word}{Style.RESET_ALL}"
            else: colorized_point += word
        
        wrapped_colorized_point = '\n'.join([colorized_point[j:j+MAX_LINE_WIDTH] for j in range(0, len(colorized_point), MAX_LINE_WIDTH)])
        display_lines.append(f"{Fore.LIGHTCYAN_EX}- {wrapped_colorized_point}{Style.RESET_ALL}")
    
    for line in display_lines:
        if len(line) > MAX_VERBIAGE_CHARS:
            line = line[:MAX_VERBIAGE_CHARS] + "..."
        print(line)
        
    summary_info = {'file_path': os.path.abspath(file_path), 'primary_category': primary_category, 'all_categories': all_valid_categories, 'summary': 

summary_text_for_file}

    if settings['file_action_choice'] != 'none' and all_valid_categories:
        for category in all_valid_categories:
            sanitized_cat = sanitize_for_path(category)
            destination_folder = os.path.join(settings['output_folder_path'], sanitized_cat)
            os.makedirs(destination_folder, exist_ok=True)
            prefix = file_management_settings.get(category, {}).get('prefix', '')
            new_file_name = f"{prefix}_{cleaned_file_name}" if prefix else cleaned_file_name
            
            destination_path = os.path.join(destination_folder, new_file_name)
            
            if settings['file_action_choice'] != 'shortcut':
                destination_path_with_ext = destination_path + os.path.splitext(file_path)[1]
            else:
                destination_path_with_ext = destination_path
            
            try:
                if settings['file_action_choice'] == 'shortcut':
                    create_shortcut(file_path, destination_path)
                    print(f"{Fore.CYAN}Created shortcut in -> {destination_folder}{Style.RESET_ALL}")
                elif category == primary_category:
                    if os.path.exists(file_path):
                        shutil.move(file_path, destination_path_with_ext)
                        print(f"{Fore.CYAN}Moved and renamed: {file_path} -> {destination_path_with_ext}{Style.RESET_ALL}")
                        if settings['tagging_enabled_choice'] in ['y', 'yes']:
                            sec_cat_names = [cat for cat, score in secondary_categories]
                            tag_file_properties(destination_path_with_ext, primary_category, sec_cat_names)
                    elif settings['file_action_choice'] == 'copy':
                        shutil.copy2(file_path, destination_path_with_ext)
                        print(f"{Fore.CYAN}Copied and renamed: {file_path} -> {destination_path_with_ext}{Style.RESET_ALL}")
                        if settings['tagging_enabled_choice'] in ['y', 'yes']:
                            sec_cat_names = [cat for cat, score in secondary_categories]
                            tag_file_properties(destination_path_with_ext, primary_category, sec_cat_names)
            except Exception as e: print(f"{Fore.LIGHTBLACK_EX}Error during file operation: {e}{Style.RESET_ALL}")
    return summary_info


# --- Main Execution ---
if __name__ == "__main__":
    multiprocessing.freeze_support()
    initialize_models()
    try:
        CATEGORIES = load_and_edit_categories()
        settings_file = "settings.json"
        defaults = {
            'pdf_chunk_length': 2048, 'generic_chunk_length': 250, 'max_summary_length': 150, 'min_summary_length': 40,
            'pdf_padding': 0, 'pdf_chunks': 1, 'generic_padding': 0, 'generic_chunks': 1,
            'confidence_threshold': 50, 'secondary_confidence_threshold': 20, 'tagging_enabled_choice': 'n',
            'add_prefix_choice': 'n', 'global_prefix_choice': 'n', 'global_prefix': '', 'special_char': '',
            'selected_file_types': ['.pdf', '.docx', '.txt', '.xlsx', '.pptx'], 'file_action_choice': 'none',
            'last_list_file': '', 'last_destination_folder': '', 'source_choice': 'f', 'preprocess_choice': 'n',
            'summarizer_timeout': 60,
            'long_pdf_threshold': 10,
            'long_text_threshold': 1000,
            'short_pdf_chunk_length': 1025,
            'long_pdf_chunk_length': 4096,
            'short_text_chunk_length': 1025,
            'long_text_chunk_length': 4096,
            'long_pdf_max_summary_length': 200,
            'long_pdf_min_summary_length': 50,
            'short_pdf_max_summary_length': 150,
            'short_pdf_min_summary_length': 40,
            'short_pdf_chunks_list': '',
            'long_pdf_chunks_list': '',
            'short_text_chunks_list': '',
            'long_text_chunks_list': '',
            'generic_chunks_list': '',
            'show_chunk_verbiage': 'n',
            'filter_keywords': '',
            'negative_filter_keywords': '',
            'keyword_context_before': 50,
            'keyword_context_after': 50,
            'culling_regex': '',
            'exclude_regex': '',
            'culling_mode': 'c'
        }

        settings_to_save = defaults.copy()
        if os.path.exists(settings_file):
            print(f"{Fore.GREEN}Found saved settings. Load them? (y/n) [default: y]: {Style.RESET_ALL}", end="")
            user_input = input().strip().lower()
            if user_input not in ['n', 'no']:
                try:
                    with open(settings_file, 'r') as f:
                        saved_settings = json.load(f)
                    settings_to_save.update(saved_settings)
                    print(f"{Fore.CYAN}Settings loaded successfully.{Style.RESET_ALL}")
                except Exception as e:
                    print(f"{Fore.RED}Error loading settings: {e}.{Style.RESET_ALL}")

        print(f"{Fore.GREEN}Do you want to see the extended help content and regex examples? (y/n) [default: no]: {Style.RESET_ALL}", end="")
        help_prompt = input().strip().lower()
        if help_prompt in ['y', 'yes']:
            print(f"\n{Fore.GREEN}# Program Overview{Style.RESET_ALL}")
            print(f"{Fore.CYAN}The core function of this program is to intelligently process documents, summarize their content, and categorize them based 

on a user-defined list of categories. The goal is to provide a clean, automated workflow for managing a large number of files. It uses two main AI models: 

a T5-small summarization model to condense text and a Zero-Shot Classification model to categorize based on provided labels. Key features include 

automatic file pre-processing, customizable chunking and summarization, and a dynamic file management system that can move, copy, or create shortcuts to 

files based on their categories. The advanced keyword and regex settings provide granular control, allowing you to fine-tune the analysis to your specific 

needs.{Style.RESET_ALL}")
            
            print(f"\n{Fore.GREEN}# Menu Settings Overview{Style.RESET_ALL}")
            print(f"{Fore.CYAN}The menu is designed to be user-friendly and persistent. Your settings are saved to 'settings.json' so you don't have to 

re-enter them for every session. You can customize core thresholds, chunking strategies for different file types, and file management actions. The 

'Keyword & Regex Settings' section gives you advanced control over what text gets processed.{Style.RESET_ALL}")

            print(f"\n{Fore.GREEN}# Advanced Regex Help: Common Patterns{Style.RESET_ALL}")
            print(f"{Fore.CYAN}Regular expressions are a powerful tool for searching and manipulating text based on patterns. Here are some useful 

patterns for culling and excluding text in your documents:{Style.RESET_ALL}\n")
            
            print(f"{Fore.GREEN}1. Email Addresses{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Matches a typical email address format. Useful for excluding personal contact information.{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{{2,}}`{Style.RESET_ALL}")
            
            print(f"{Fore.GREEN}\n2. Phone Numbers{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Matches common US phone number formats, with or without area codes and separators.{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `\d{{3}}[-\s]?\d{{3}}[-\s]?\d{{4}}`{Style.RESET_ALL}")
            
            print(f"{Fore.GREEN}\n3. Dates{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Matches dates in the format MM/DD/YYYY or MM-DD-YYYY.{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `\d{{2}}[/.-]\d{{2}}[/.-]\d{{4}}`{Style.RESET_ALL}")
            
            print(f"{Fore.GREEN}\n4. URLs/Links{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Matches common web addresses, starting with http:// or https://.{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `https?:\/\/[^\s]+`{Style.RESET_ALL}")
            
            print(f"{Fore.GREEN}\n5. Social Security Numbers{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Matches the standard SSN format (###-##-####).{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `\d{{3}}-\d{{2}}-\d{{4}}`{Style.RESET_ALL}")
            
            print(f"{Fore.GREEN}\n6. Financial Figures{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Matches dollar amounts, including optional commas and decimals.{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `\$\d{{1,3}}(?:,?\d{{3}})*(?:\.\d{{2}})?`{Style.RESET_ALL}")
            
            print(f"{Fore.GREEN}\n7. All Punctuation and Symbols{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Matches most standard punctuation. Useful for cleaning text.{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `[^\w\s]`{Style.RESET_ALL}")
            
            print(f"{Fore.GREEN}\n8. Specific Culling Keywords{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Combines multiple words with the OR operator. Matches any of the words as a whole word.{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `\b(keyword1|keyword2|anotherkeyword)\b`{Style.RESET_ALL}")
            
            print(f"{Fore.GREEN}\n9. Removing Headers/Footers{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Matches common header/footer text patterns at the beginning or end of a page. (e.g., 'Page X of Y').{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `^Page\s\d+\sof\s\d+$`{Style.RESET_ALL}")
            
            print(f"{Fore.GREEN}\n10. Finding Two-Word Phrases{Style.RESET_ALL}")
            print(f"{Fore.CYAN}  Matches two consecutive words, which can be useful for finding phrases.{Style.RESET_ALL}")
            print(fr"{Fore.YELLOW}  `\b(\w+\s\w+)\b`{Style.RESET_ALL}")
            
            print(f"\n{Fore.GREEN}---{Style.RESET_ALL}")

        print(f"{Fore.GREEN}Process files from a [f]older or a [l]ist file? [default: {settings_to_save['source_choice']}]: {Style.RESET_ALL}", end="")
        source_choice = input().strip().lower() or settings_to_save['source_choice']

        initial_file_list, output_folder_path, list_file_path = "", "", settings_to_save.get('last_list_file', '')
        last_destination_folder = settings_to_save.get('last_destination_folder', '')

        if source_choice == 'f':
            print(f"{Fore.GREEN}Enter the folder path to scan: {Style.RESET_ALL}", end="")
            folder_path = input().strip()
            print(f"{Fore.GREEN}Send organized files to [s]ubdirectory or [d]efined location?: {Style.RESET_ALL}", end="")
            dest_choice = input().strip().lower()
            if dest_choice == 's':
                output_folder_path = folder_path
            else:
                print(f"{Fore.GREEN}Enter the destination folder path [default: {last_destination_folder}]: {Style.RESET_ALL}", end="")
                output_folder_path = input().strip() or last_destination_folder
            print(f"{Fore.GREEN}Scan subdirectories of the source folder? (y/n) [default: n]: {Style.RESET_ALL}", end="")
            scan_sub = input().strip().lower() in ['y', 'yes']
            
            initial_file_list = []
            walk_target = os.walk(folder_path) if scan_sub else [(os.path.dirname(folder_path) if folder_path else '.', [], os.listdir(folder_path or 

'.'))]
            for root, _, files in walk_target:
                for file in files:
                    if os.path.splitext(file)[1].lower() in settings_to_save['selected_file_types']: initial_file_list.append(os.path.join(root, file))
        elif source_choice == 'l':
            print(f"{Fore.GREEN}Enter path to the list file [default: {settings_to_save['last_list_file']}]: {Style.RESET_ALL}", end="")
            list_file_path = input().strip() or settings_to_save['last_list_file']
            print(f"{Fore.GREEN}Enter a destination folder path [default: {last_destination_folder}]: {Style.RESET_ALL}", end="")
            output_folder_path = input().strip() or last_destination_folder

            initial_file_list = []
            with open(list_file_path, 'r', encoding='utf-8') as f:
                for line in f:
                    file_path = line.strip()
                    if os.path.isfile(file_path) and os.path.splitext(file_path)[1].lower() in settings_to_save['selected_file_types']: 

initial_file_list.append(file_path)

        if not os.path.isdir(output_folder_path): os.makedirs(output_folder_path, exist_ok=True)
        settings_to_save['output_folder_path'] = output_folder_path

        action_map = {'m': 'move', 'c': 'copy', 's': 'shortcut', 'd': 'none'}
        print(f"{Fore.GREEN}Action: [M]ove, [C]opy, create [S]hortcut, or [D]o nothing? [default: {settings_to_save['file_action_choice']}]: 

{Style.RESET_ALL}", end="")
        action_input = input().strip().lower()
        settings_to_save['file_action_choice'] = action_map.get(action_input, settings_to_save['file_action_choice'] if not action_input else 'none')
        print(f"{Fore.GREEN}Pre-process by filename? (y/n) [default: {settings_to_save['preprocess_choice']}]: {Style.RESET_ALL}", end="")
        settings_to_save['preprocess_choice'] = input().strip().lower() or settings_to_save['preprocess_choice']

        print(f"\n{Fore.GREEN}--- Configuration Settings ---{Style.RESET_ALL}")
        
        print(f"\n{Fore.GREEN}# Core Settings{Style.RESET_ALL}")
        print(f"{Fore.CYAN}# These settings control the overall behavior of the summarizer and categorizer.{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Enter summarizer timeout in seconds [default: {settings_to_save['summarizer_timeout']}]: {Style.RESET_ALL}", end="")
        settings_to_save['summarizer_timeout'] = int(input().strip() or settings_to_save['summarizer_timeout'])
        print(f"{Fore.GREEN}Enter primary confidence threshold (0-100) [default: {settings_to_save['confidence_threshold']}]: {Style.RESET_ALL}", end="")
        settings_to_save['confidence_threshold'] = int(input().strip() or settings_to_save['confidence_threshold'])
        print(f"{Fore.GREEN}Enter secondary confidence threshold (0-100) [default: {settings_to_save['secondary_confidence_threshold']}]: 

{Style.RESET_ALL}", end="")
        settings_to_save['secondary_confidence_threshold'] = int(input().strip() or settings_to_save['secondary_confidence_threshold'])
        print(f"{Fore.GREEN}Tag files with category keywords? (y/n) [default: {settings_to_save['tagging_enabled_choice']}]: {Style.RESET_ALL}", end="")
        settings_to_save['tagging_enabled_choice'] = input().strip().lower() or settings_to_save['tagging_enabled_choice']
        
        print(f"\n{Fore.GREEN}---{Style.RESET_ALL}")
        
        print(f"\n{Fore.GREEN}# Chunk & Text Settings{Style.RESET_ALL}")
        print(f"{Fore.CYAN}# These settings are used for parsing and preparing text from different file types.{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Enter non-PDF chunk length [default: {settings_to_save['generic_chunk_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['generic_chunk_length'] = int(input().strip() or settings_to_save['generic_chunk_length'])
        print(f"{Fore.GREEN}Enter SHORT PDF chunk length [default: {settings_to_save.get('short_pdf_chunk_length', defaults['short_pdf_chunk_length'])}]: 

{Style.RESET_ALL}", end="")
        settings_to_save['short_pdf_chunk_length'] = int(input().strip() or settings_to_save.get('short_pdf_chunk_length', defaults

['short_pdf_chunk_length']))
        print(f"{Fore.GREEN}Enter LONG PDF chunk length [default: {settings_to_save.get('long_pdf_chunk_length', defaults['long_pdf_chunk_length'])}]: 

{Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_chunk_length'] = int(input().strip() or settings_to_save.get('long_pdf_chunk_length', defaults

['long_pdf_chunk_length']))
        
        print(f"{Fore.GREEN}Enter SHORT PDF max summary length [default: {settings_to_save['max_summary_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['max_summary_length'] = int(input().strip() or settings_to_save['max_summary_length'])
        print(f"{Fore.GREEN}Enter SHORT PDF min summary length [default: {settings_to_save['min_summary_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['min_summary_length'] = int(input().strip() or settings_to_save['min_summary_length'])
        print(f"{Fore.GREEN}Enter the LONG PDF threshold (in pages) [default: {settings_to_save['long_pdf_threshold']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_threshold'] = int(input().strip() or settings_to_save['long_pdf_threshold'])
        print(f"{Fore.GREEN}Enter LONG PDF max summary length [default: {settings_to_save['long_pdf_max_summary_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_max_summary_length'] = int(input().strip() or settings_to_save['long_pdf_max_summary_length'])
        print(f"{Fore.GREEN}Enter LONG PDF min summary length [default: {settings_to_save['long_pdf_min_summary_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_min_summary_length'] = int(input().strip() or settings_to_save['long_pdf_min_summary_length'])
        
        print(f"\n{Fore.GREEN}---{Style.RESET_ALL}")

        print(f"\n{Fore.GREEN}# Chunk Selection (comma-separated list/ranges){Style.RESET_ALL}")
        print(f"{Fore.CYAN}# These settings define which chunks of a document are read and processed.{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Enter chunks for SHORT PDFs (e.g., '1, 3-5') [default: {settings_to_save.get('short_pdf_chunks_list', '')}]: 

{Style.RESET_ALL}", end="")
        settings_to_save['short_pdf_chunks_list'] = input().strip() or settings_to_save.get('short_pdf_chunks_list', '')
        print(f"{Fore.GREEN}Enter chunks for LONG PDFs (e.g., '1-2, 10') [default: {settings_to_save.get('long_pdf_chunks_list', '')}]: 

{Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_chunks_list'] = input().strip() or settings_to_save.get('long_pdf_chunks_list', '')
        print(f"{Fore.GREEN}Enter chunks for SHORT text files (e.g., '1') [default: {settings_to_save.get('short_text_chunks_list', '')}]: 

{Style.RESET_ALL}", end="")
        settings_to_save['short_text_chunks_list'] = input().strip() or settings_to_save.get('short_text_chunks_list', '')
        print(f"{Fore.GREEN}Enter chunks for LONG text files (e.g., '1-5') [default: {settings_to_save.get('long_text_chunks_list', '')}]: 

{Style.RESET_ALL}", end="")
        settings_to_save['long_text_chunks_list'] = input().strip() or settings_to_save.get('long_text_chunks_list', '')
        print(f"{Fore.GREEN}Enter the LONG TEXT threshold (in words) [default: {settings_to_save['long_text_threshold']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_text_threshold'] = int(input().strip() or settings_to_save['long_text_threshold'])
        print(f"{Fore.GREEN}Enter SHORT text chunk length [default: {settings_to_save['short_text_chunk_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['short_text_chunk_length'] = int(input().strip() or settings_to_save['short_text_chunk_length'])
        print(f"{Fore.GREEN}Enter LONG text chunk length [default: {settings_to_save['long_text_chunk_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_text_chunk_length'] = int(input().strip() or settings_to_save['long_text_chunk_length'])
        
        print(f"\n{Fore.GREEN}---{Style.RESET_ALL}")
        
        print(f"\n{Fore.GREEN}# Keyword & Regex Settings{Style.RESET_ALL}")
        print(f"{Fore.CYAN}# These settings provide fine-grained control over keyword filtering and text culling. You can use these to create a more 

focused and relevant input for the summarizer.{Style.RESET_ALL}")
        
        current_show_verbiage = settings_to_save.get('show_chunk_verbiage', 'n')
        print(f"{Fore.GREEN}Show extraction and summary text for each chunk? (y/n) [default: {current_show_verbiage}]: {Style.RESET_ALL}", end="")
        settings_to_save['show_chunk_verbiage'] = input().strip().lower() or current_show_verbiage
        
        print(f"\n{Fore.CYAN}# Culling Mode: Choose your method for pre-processing text for summarization. Fall back will be to defined chunk size.

{Style.RESET_ALL}")
        print(f"{Fore.CYAN}#      'c' - Cull to complete sentences (Default).{Style.RESET_ALL}")
        print(f"{Fore.CYAN}#      's' - Extract the complete sentence(s) containing the keyword.{Style.RESET_ALL}")
        print(f"{Fore.CYAN}#      'b' - Extract the sentence before, the sentence containing, and the sentence after the keyword.{Style.RESET_ALL}")
        print(f"{Fore.CYAN}#      '2' - Extract two sentences before, the sentence containing, and two sentences after the keyword.{Style.RESET_ALL}")
        print(f"{Fore.CYAN}#      'p' - Extract the full paragraph containing the keyword.{Style.RESET_ALL}")
        print(f"{Fore.CYAN}#      'r' - Use a custom regex string to define the culling pattern.{Style.RESET_ALL}")
        print(f"{Fore.CYAN}#      'n' - No culling; summarize the entire extracted text.{Style.RESET_ALL}")

        culling_mode_options = ['c', 's', 'b', '2', 'p', 'r', 'n']
        print(f"{Fore.GREEN}Select a culling mode (c/s/b/2/p/r/n) [default: {settings_to_save.get('culling_mode', 'c')}]: {Style.RESET_ALL}", end="")
        culling_mode_choice = input().strip().lower()
        if culling_mode_choice and culling_mode_choice in culling_mode_options:
            settings_to_save['culling_mode'] = culling_mode_choice
        else:
            settings_to_save['culling_mode'] = settings_to_save.get('culling_mode', 'c')
        
        if settings_to_save['culling_mode'] in ['s', 'b', '2', 'p', 'k']:
            print(f"\n{Fore.CYAN}# These keywords are used to locate the relevant sentences/paragraphs for culling.{Style.RESET_ALL}")
            current_filter_keywords = settings_to_save.get('filter_keywords', '')
            print(f"{Fore.GREEN}Enter filter keywords (comma-separated) [default: '{current_filter_keywords}']: {Style.RESET_ALL}", end="")
            settings_to_save['filter_keywords'] = input().strip() or current_filter_keywords
        
        if settings_to_save['culling_mode'] == 'k':
            print(f"\n{Fore.CYAN}# When culling by keywords, these settings define the scope.{Style.RESET_ALL}")
            current_filter_keywords = settings_to_save.get('filter_keywords', '')
            print(f"{Fore.GREEN}Enter filter keywords (comma-separated) [default: '{current_filter_keywords}']: {Style.RESET_ALL}", end="")
            settings_to_save['filter_keywords'] = input().strip() or current_filter_keywords
            
            current_context_before = settings_to_save.get('keyword_context_before', 50)
            print(f"{Fore.GREEN}Enter characters to include BEFORE a keyword [default: {current_context_before}]: {Style.RESET_ALL}", end="")
            settings_to_save['keyword_context_before'] = int(input().strip() or current_context_before)
            
            current_context_after = settings_to_save.get('keyword_context_after', 50)
            print(f"{Fore.GREEN}Enter characters to include AFTER a keyword [default: {current_context_after}]: {Style.RESET_ALL}", end="")
            settings_to_save['keyword_context_after'] = int(input().strip() or current_context_after)
            
        elif settings_to_save['culling_mode'] == 'r':
            print(f"\n{Fore.CYAN}# When culling by regex, provide your custom pattern.{Style.RESET_ALL}")
            current_culling_regex = settings_to_save.get('culling_regex', '')
            print(f"{Fore.GREEN}Enter regex for culling text [default: '{current_culling_regex}']: {Style.RESET_ALL}", end="")
            settings_to_save['culling_regex'] = input().strip() or current_culling_regex
            
        print(f"\n{Fore.CYAN}# These keywords are used to skip processing files that are likely irrelevant to your interests. Files containing any of 

these keywords will be skipped.{Style.RESET_ALL}")
        current_negative_keywords = settings_to_save.get('negative_filter_keywords', '')
        print(f"{Fore.GREEN}Enter keywords to filter OUT (comma-separated) [default: '{current_negative_keywords}']: {Style.RESET_ALL}", end="")
        settings_to_save['negative_filter_keywords'] = input().strip() or current_negative_keywords

        print(f"\n{Fore.CYAN}# This regex is used to exclude files from processing. It will be applied to the entire file content. Files matching this 

pattern will be skipped entirely.{Style.RESET_ALL}")
        current_exclude_regex = settings_to_save.get('exclude_regex', '')
        print(f"{Fore.GREEN}Enter regex to exclude files [default: '{current_exclude_regex}']: {Style.RESET_ALL}", end="")
        settings_to_save['exclude_regex'] = input().strip() or current_exclude_regex
        
        print(f"\n{Fore.GREEN}---{Style.RESET_ALL}")
        
        if settings_to_save['file_action_choice'] != 'none':
            print(f"{Fore.GREEN}Add a prefix to organized files? (y/n) [default: {settings_to_save['add_prefix_choice']}]: {Style.RESET_ALL}", end="")
            user_input = input().strip().lower()
            settings_to_save['add_prefix_choice'] = user_input or settings_to_save['add_prefix_choice']
            if settings_to_save['add_prefix_choice'] in ['y', 'yes']:
                print(f"{Fore.GREEN}Apply a global prefix? (y/n) [default: {settings_to_save['global_prefix_choice']}]: {Style.RESET_ALL}", end="")
                user_input = input().strip().lower()
                settings_to_save['global_prefix_choice'] = user_input or settings_to_save['global_prefix_choice']
                if settings_to_save['global_prefix_choice'] in ['y', 'yes']:
                    print(f"{Fore.GREEN}Enter a global prefix (max 12 chars) [default: {settings_to_save['global_prefix']}]: {Style.RESET_ALL}", end="")
                    settings_to_save['global_prefix'] = input().strip() or settings_to_save['global_prefix']
                    print(f"{Fore.GREEN}Enter a special character (1-3 chars) [default: {settings_to_save['special_char']}]: {Style.RESET_ALL}", end="")
                    settings_to_save['special_char'] = input().strip() or settings_to_save['special_char']

        settings_to_save['source_choice'] = source_choice
        settings_to_save['last_destination_folder'] = output_folder_path
        if source_choice == 'l': settings_to_save['last_list_file'] = list_file_path
        with open(settings_file, 'w') as f: json.dump(settings_to_save, f, indent=4)
        print(f"\n{Fore.CYAN}Settings saved. Starting process...{Style.RESET_ALL}")

        file_management_settings = {}
        if settings_to_save['file_action_choice'] != 'none' and output_folder_path:
            for category in CATEGORIES:
                if category != "Other":
                    prefix_str = ""
                    if settings_to_save['add_prefix_choice'] in ['y', 'yes']:
                        if settings_to_save['global_prefix_choice'] in ['y', 'yes']: prefix_str = f"{settings_to_save['global_prefix']}{settings_to_save

['special_char']}"
                        else: prefix_str = sanitize_for_path(category).replace(' ', '')[:12].upper()
                    file_management_settings[category] = {'prefix': prefix_str, 'destination': os.path.join(output_folder_path, sanitize_for_path

(category))}

        # Pre-build a cache of all existing processed files to speed up subsequent runs
        existing_files_cache = set()
        if settings_to_save['file_action_choice'] != 'none' and output_folder_path:
            print(f"{Fore.CYAN}Step 2: Building cache of existing files for faster re-runs...{Style.RESET_ALL}")
            # Add all user-defined categories
            for category in CATEGORIES:
                sanitized_cat = sanitize_for_path(category)
                dest_folder = os.path.join(output_folder_path, sanitized_cat)
                if os.path.isdir(dest_folder):
                    for fname in os.listdir(dest_folder):
                        existing_files_cache.add(fname)
            
            # Add special categories
            for special_cat in ["contains no text", "unreadable file"]:
                sanitized_cat = sanitize_for_path(special_cat)
                dest_folder = os.path.join(output_folder_path, sanitized_cat)
                if os.path.isdir(dest_folder):
                    for fname in os.listdir(dest_folder):
                        existing_files_cache.add(fname)
            print(f"{Fore.CYAN}Cache built with {len(existing_files_cache)} entries.{Style.RESET_ALL}")


        files_to_process = initial_file_list
        if settings_to_save['preprocess_choice'] in ['y', 'yes'] and settings_to_save['file_action_choice'] != 'none':
            files_to_process = preprocess_by_filename(initial_file_list, CATEGORIES, settings_to_save['file_action_choice'], output_folder_path, 

file_management_settings, settings_to_save)
        
        # New pre-filtering step
        if settings_to_save['file_action_choice'] != 'none':
            print(f"{Fore.CYAN}Step 3: Pre-filtering files against cache...{Style.RESET_ALL}")
            filtered_files = []
            total_initial_files = len(files_to_process)
            for file_path in files_to_process:
                should_skip = False
                
                file_name_no_ext = os.path.splitext(os.path.basename(file_path))[0]
                cleaned_file_name = file_name_no_ext.replace('_', ' ').replace('-', ' ')
                
                # Check against user-defined categories
                for category in CATEGORIES:
                    if category == "Other": continue
                    prefix = file_management_settings.get(category, {}).get('prefix', '')
                    new_file_name_no_ext = f"{prefix}_{cleaned_file_name}" if prefix else cleaned_file_name
                    
                    if settings_to_save['file_action_choice'] == 'shortcut':
                        expected_output_name = new_file_name_no_ext + ".url" if sys.platform == "win32" else new_file_name_no_ext
                    else:
                        expected_output_name = new_file_name_no_ext + os.path.splitext(file_path)[1]

                    if expected_output_name in existing_files_cache:
                        print(f"{Fore.LIGHTBLACK_EX}Skipping pre-filtered: '{os.path.basename(file_path)}' found in '{sanitize_for_path(category)}' 

subfolder.{Style.RESET_ALL}")
                        should_skip = True
                        break
                        
                # Check against special categories
                for special_cat in ["contains no text", "unreadable file"]:
                    sanitized_cat = sanitize_for_path(special_cat)
                    if settings_to_save['file_action_choice'] == 'shortcut':
                        expected_output_name = cleaned_file_name + ".url" if sys.platform == "win32" else cleaned_file_name
                    else:
                        expected_output_name = cleaned_file_name + os.path.splitext(file_path)[1]

                    if expected_output_name in existing_files_cache:
                        print(f"{Fore.LIGHTBLACK_EX}Skipping pre-filtered: '{os.path.basename(file_path)}' found in '{sanitized_cat}' subfolder.

{Style.RESET_ALL}")
                        should_skip = True
                        break

                if not should_skip:
                    filtered_files.append(file_path)
            
            files_to_process = filtered_files
            
            print(f"{Fore.CYAN}Pre-filtering complete: {len(files_to_process)} new files to process out of a total of {total_initial_files} from the 

initial list.{Style.RESET_ALL}")


        all_summaries = []
        total_files = len(files_to_process)
        for i, file_to_process in enumerate(files_to_process):
            summary_info = process_single_file((file_to_process, i, total_files, CATEGORIES, file_management_settings, settings_to_save))
            if summary_info:
                all_summaries.append(summary_info)

        if all_summaries and output_folder_path:
            summary_file_path = os.path.join(output_folder_path, "all_summaries.txt")
            with open(summary_file_path, "w", encoding="utf-8") as f:
                for summary in all_summaries:
                    f.write("="*50 + "\n")
                    f.write(f"File: {summary['file_path']}\n")
                    f.write(f"Primary Category: {summary['primary_category']}\n")
                    if len(summary['all_categories']) > 1: f.write(f"All Categories: {', '.join(summary['all_categories'])}\n")
                    f.write(f"\nSummary:\n{summary['summary']}\n")
            print(f"Summaries saved to: {summary_file_path}")

    except KeyboardInterrupt:
        print(f"{Fore.LIGHTBLACK_EX}\nProcess interrupted.{Style.RESET_ALL}")
    except Exception as e:
        print(f"{Fore.RED}\nAn unexpected error occurred: {e}{Style.RESET_ALL}")
