# comboaisort25.py
import os
import fitz  # PyMuPDF for PDF
from docx import Document
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook
from transformers import pipeline
import torch
from transformers import AutoTokenizer, AutoModelForSequenceClassification
import shutil
import json
import colorama
from colorama import Fore, Style
import re
import sys

# Initialize colorama to auto-reset styles after each print
colorama.init(autoreset=True)

# Check for GPU availability and set the device
if torch.cuda.is_available():
    device = "cuda"
    print("Device set to use GPU.")
else:
    device = "cpu"
    print("GPU not found. Device set to use CPU.")

# --- Configuration ---
# These are now set via user input or loaded from a file.

# Define the categories for categorization
CATEGORIES = [
    "Health",
    "History",
    "Weather",
    "do not combine general categories with categories containing targeted things within it (delete this)"
]

# --- Load and Edit Categories Function ---
def load_and_edit_categories():
    """
    Loads categories from a file, prompts the user to edit them, and saves the final list.
    """
    category_file = "categories.json"
    default_categories = [
        "Health",
        "History",
        "Weather",
        "do not combine general categories with categories containing targeted things within it (delete this)"
    ]
    categories = default_categories
    
    try:
        if os.path.exists(category_file):
            with open(category_file, 'r', encoding='utf-8') as f:
                categories = json.load(f)
            print("Step 1.2: Loaded categories from previous session.")
    except Exception as e:
        print(f"Error loading categories file: {e}. Using default categories.")

    print("\n--- Current Categories ---")
    for i, cat in enumerate(categories):
        print(f"[{i+1}] {cat}")
    print("--------------------------")
    
    edit_choice = input("Do you want to edit these categories? (y/n) [default: no]: ").strip().lower()
    if edit_choice in ['y', 'yes']:
        instructions = "\nEditing categories. Enter 'add <new_category>', 'remove <number>', 'edit <number> <new_category>', 'list', or 'done' to finish."
        print(instructions)
        while True:
            command = input("> ").strip().lower()
            if command == 'done':
                break
            
            parts = command.split()
            if not parts:
                continue

            action = parts[0]
            if action == 'add' and len(parts) >= 2:
                new_cat = " ".join(parts[1:])
                categories.append(new_cat)
                print(f"Added: {new_cat}")
                print("\n--- Current Categories ---")
                for i, cat in enumerate(categories):
                    print(f"[{i+1}] {cat}")
                print("--------------------------")
                print(instructions)
            elif action == 'remove' and len(parts) == 2 and parts[1].isdigit():
                idx = int(parts[1]) - 1
                if 0 <= idx < len(categories):
                    removed_cat = categories.pop(idx)
                    print(f"Removed: {removed_cat}")
                    print("\n--- Current Categories ---")
                    for i, cat in enumerate(categories):
                        print(f"[{i+1}] {cat}")
                    print("--------------------------")
                    print(instructions)
                else:
                    print("Invalid category number.")
            elif action == 'edit' and len(parts) >= 3 and parts[1].isdigit():
                idx = int(parts[1]) - 1
                if 0 <= idx < len(categories):
                    new_cat = " ".join(parts[2:])
                    old_cat = categories[idx]
                    categories[idx] = new_cat
                    print(f"Edited: '{old_cat}' to '{new_cat}'")
                    print("\n--- Current Categories ---")
                    for i, cat in enumerate(categories):
                        print(f"[{i+1}] {cat}")
                    print("--------------------------")
                    print(instructions)
                else:
                    print("Invalid category number.")
            elif action == 'list':
                print("\n--- Current Categories ---")
                for i, cat in enumerate(categories):
                    print(f"[{i+1}] {cat}")
                print("--------------------------")
                print(instructions)
            else:
                print("Invalid command. Please try again.")
    
    # Save the final list of categories
    with open(category_file, 'w', encoding='utf-8') as f:
        json.dump(categories, f, indent=4)
    print("\nStep 1.3: Categories saved for next session.")
    
    return categories

# --- Load Models ---
print("\nStep 1: Loading summarization and classification models...")
print("If this is the first time you are running the script, a large file download will begin now. Please wait for it to complete.")

# Load the SMALLER summarization pipeline for speed
summarizer = pipeline(
    "summarization",
    model="t5-small", 
    max_length=1024, 
    device=device,
    framework="pt"  # Explicitly tell the pipeline to use PyTorch
)

# Load a dedicated model for text classification
try:
    classifier = pipeline(
        "zero-shot-classification",
        model="MoritzLaurer/xtremedistil-l6-h256-zeroshot-v1.1-all-33",
        device=device
    )
    print("Step 1.1: Classification model loaded successfully.")
except Exception as e:
    print(f"Error loading classification model: {e}")
    classifier = None
    print("Step 1.1: Falling back to 'Other' for all classifications.")


# --- File Reading Functions ---
def read_chunks_from_file(file_path, start_chunk, end_chunk, max_chunk_length):
    """
    Efficiently reads a specific range of word chunks from any file type.
    It stops reading the file once the desired number of chunks has been collected.
    """
    ext = os.path.splitext(file_path)[1].lower()
    words = []
    total_words_needed = end_chunk * max_chunk_length

    try:
        if ext == ".pdf":
            with fitz.open(file_path) as doc:
                for page in doc:
                    words.extend(page.get_text().split())
                    if len(words) >= total_words_needed:
                        break
        elif ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                words.extend(para.text.split())
                if len(words) >= total_words_needed:
                    break
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                # This is less efficient for huge txt files, but is the standard way
                words = f.read().split()
        elif ext == ".xlsx":
            sheets = pd.ExcelFile(file_path).sheet_names
            for sheet in sheets:
                df = pd.read_excel(file_path, sheet_name=sheet)
                words.extend(df.to_string(index=False).split())
                if len(words) >= total_words_needed:
                    break
        elif ext == ".pptx":
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        words.extend(shape.text.split())
                        if len(words) >= total_words_needed:
                            break
                if len(words) >= total_words_needed:
                    break
    except Exception as e:
        print(f"Error reading {ext} file: {e}")
        return ""

    start_index = (start_chunk - 1) * max_chunk_length
    end_index = end_chunk * max_chunk_length
    
    # Return the requested slice of words, joined back into a string
    return " ".join(words[start_index:end_index])

# --- Summarization Function ---
def summarize_text(text, max_chunk_length, max_summary_length, min_summary_length):
    """
    Generates a summary from the provided text by processing it in word chunks.
    """
    if not text.strip():
        return ["No text to summarize."]

    words = text.split()
    chunks = [" ".join(words[i:i + max_chunk_length]) for i in range(0, len(words), max_chunk_length)]
    
    summaries = []
    
    for i, chunk in enumerate(chunks):
        print(f"Step 7.1.1: Summarizing chunk {i + 1} of {len(chunks)}...")
        try:
            summary = summarizer(
                chunk, 
                max_length=max_summary_length, 
                min_length=min_summary_length, 
                truncation=True
            )
            summaries.append(summary[0]['summary_text'])
            print(f"Step 7.1.2: Summarization for chunk {i + 1} complete.")
        except Exception as e:
            print(f"Step 7.1.3: Error during summarization of chunk {i + 1}: {e}")
            summaries.append("Summary failed for this chunk.")
    
    combined_summary = ". ".join(summaries)
    return combined_summary.split(". ")
    
# --- Categorization Function ---
def categorize_summary(summary_text, categories):
    """
    Categorizes a summary using a zero-shot classification model.
    Returns the full results dictionary from the classifier.
    """
    if classifier is None:
        return {'labels': ['Other'], 'scores': [1.0]}

    print("Step 8.1: Sending summary to classifier model...")
    try:
        results = classifier(summary_text, candidate_labels=categories)
        print("Step 8.2: Categorization complete.")
        return results
    except Exception as e:
        print(f"Step 8.3: Error during AI categorization: {e}")
        return {'labels': ['Other'], 'scores': [1.0]}

# --- File Tagging Function ---
def tag_file_properties(file_path, category):
    """
    Adds the category to the file's metadata properties (keywords/tags).
    """
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if not os.path.exists(file_path):
            print(f"Tagging skipped: File not found at {file_path}")
            return

        if ext == ".pdf":
            doc = fitz.open(file_path)
            metadata = doc.metadata
            keywords = metadata.get("keywords", "")
            if category not in keywords.split(';'):
                new_keywords = f"{keywords}; {category}".strip('; ')
                metadata["keywords"] = new_keywords
                doc.set_metadata(metadata)
                doc.save(file_path, incremental=True, encryption=fitz.PDF_ENCRYPT_KEEP)
            doc.close()
            print(f"Tagged '{os.path.basename(file_path)}' with category '{category}'.")
        elif ext == ".docx":
            doc = Document(file_path)
            keywords = doc.core_properties.keywords or ""
            if category not in keywords.split(';'):
                doc.core_properties.keywords = f"{keywords}; {category}".strip('; ')
                doc.save(file_path)
            print(f"Tagged '{os.path.basename(file_path)}' with category '{category}'.")
        elif ext == ".pptx":
            prs = Presentation(file_path)
            keywords = prs.core_properties.keywords or ""
            if category not in [kw.strip() for kw in keywords.split(';')]:
                prs.core_properties.keywords = f"{keywords}; {category}".strip('; ')
                prs.save(file_path)
            print(f"Tagged '{os.path.basename(file_path)}' with category '{category}'.")
        elif ext == ".xlsx":
            wb = load_workbook(file_path)
            keywords = wb.properties.keywords or ""
            if category not in keywords.split(';'):
                wb.properties.keywords = f"{keywords}; {category}".strip('; ')
                wb.save(file_path)
            print(f"Tagged '{os.path.basename(file_path)}' with category '{category}'.")
        else:
            print(f"Tagging not supported for {ext} files.")
    except Exception as e:
        print(f"Error tagging file {os.path.basename(file_path)}: {e}")

# --- Shortcut Creation Function ---
def create_shortcut(source_path, shortcut_path):
    """
    Creates a shortcut for a file, handling different OSs.
    """
    if sys.platform == "win32":
        # Create a .url file for Windows
        shortcut_path += ".url"
        with open(shortcut_path, "w") as f:
            f.write("[InternetShortcut]\n")
            f.write(f"URL=file:///{os.path.abspath(source_path)}\n")
    else:
        # Create a symbolic link for macOS/Linux
        if os.path.lexists(shortcut_path):
            os.remove(shortcut_path)
        os.symlink(source_path, shortcut_path)

# --- Main Processing Logic ---
def process_files_in_folder(file_list, categories, pdf_padding, pdf_chunks, generic_padding, generic_chunks, file_management_settings, 

confidence_threshold, pdf_chunk_length, generic_chunk_length, max_summary_length, min_summary_length, tagging_enabled, tagging_threshold, file_action, 

output_folder_path):
    """
    Walks a folder, processes supported files, and generates summaries.
    Includes detailed progress checks.
    """
    if not file_list:
        print("\nNo files to process.")
        return []

    print(f"\nStep 4: Starting main processing for {len(file_list)} files...")

    all_summaries = []

    for i, file_path in enumerate(file_list):
        print(f"\n--- Processing file {i + 1}/{len(file_list)}: {os.path.basename(file_path)} ---")
        
        # --- Check if output file already exists ---
        if file_action != 'none':
            original_basename_no_ext = os.path.splitext(os.path.basename(file_path))[0]
            
            if file_action == 'shortcut':
                expected_output_name = original_basename_no_ext + ".url" if sys.platform == "win32" else original_basename_no_ext
            else:
                expected_output_name = os.path.basename(file_path)

            already_processed = False
            for cat in categories:
                potential_dest_path = os.path.join(output_folder_path, cat, expected_output_name)
                if os.path.exists(potential_dest_path):
                    print(f"Skipping: An output file '{expected_output_name}' already exists in a category subfolder '{cat}'.")
                    already_processed = True
                    break
            if already_processed:
                continue

        is_pdf = os.path.splitext(file_path)[1].lower() == '.pdf'
        
        if is_pdf:
            padding = pdf_padding
            chunks_to_process = pdf_chunks
            max_chunk_length = pdf_chunk_length
        else:
            padding = generic_padding
            chunks_to_process = generic_chunks
            max_chunk_length = generic_chunk_length
            
        start_chunk = padding + 1
        end_chunk = padding + chunks_to_process
        
        print(f"\nStep 6: Extracting from chunk(s) {start_chunk} to {end_chunk}...")
        text = read_chunks_from_file(file_path, start_chunk, end_chunk, max_chunk_length)

        if text.strip():
            print("Step 7: Text extracted successfully. Starting summarization...")
            bullet_points = summarize_text(text, max_chunk_length, max_summary_length, min_summary_length)
            
            if bullet_points and bullet_points != ["No text to summarize."]:
                file_name_for_analysis = os.path.basename(file_path)
                file_name_no_ext = os.path.splitext(file_name_for_analysis)[0]
                cleaned_file_name = file_name_no_ext.replace('_', ' ').replace('-', ' ')
                
                all_chunks_for_classifier = [cleaned_file_name] + bullet_points
                text_for_classifier = ". ".join(all_chunks_for_classifier)
                
                print("Step 8: Categorizing summary offline...")
                results = categorize_summary(text_for_classifier, categories)
                
                top_score = results['scores'][0]
                
                if top_score * 100 >= confidence_threshold:
                    top_category = results['labels'][0]
                    other_categories = results['labels'][1:]
                else:
                    top_category = "Other"
                    other_categories = results['labels']
                    print(f"Top category '{results['labels'][0]}' with score {top_score*100:.2f}% is below threshold of {confidence_threshold}%. Assigning 

to 'Other'.")

                print("Step 9: Final summary complete.")
                
                # Tag the original file first if tagging is enabled
                if tagging_enabled and top_score * 100 >= tagging_threshold:
                    tag_file_properties(file_path, top_category)
                
                print(f"\nFile: {file_path}")
                print("Summary:\n")
                
                print(f"Category: {Fore.GREEN}{Style.BRIGHT}{top_category} ({top_score*100:.2f}%)")
                
                summary_text_for_file = ""
                for point in bullet_points:
                    colorized_point = ""
                    words_and_delimiters = re.split(r'([ ,.;:!?])', point)
                    
                    for word in words_and_delimiters:
                        if word.lower() == top_category.lower():
                            colorized_point += f"{Fore.GREEN}{Style.BRIGHT}{word}{Style.RESET_ALL}"
                        elif word.lower() in [cat.lower() for cat in other_categories]:
                            colorized_point += f"{Fore.YELLOW}{word}{Style.RESET_ALL}"
                        else:
                            colorized_point += word
                    
                    print(f"- {colorized_point.strip()}")
                    summary_text_for_file += f"- {point.strip()}\n"
                
                colorized_filename = ""
                filename_words_and_delimiters = re.split(r'([ ,.;:!?])', cleaned_file_name)
                for word in filename_words_and_delimiters:
                    if word.lower() == top_category.lower():
                        colorized_filename += f"{Fore.GREEN}{Style.BRIGHT}{word}{Style.RESET_ALL}"
                    elif word.lower() in [cat.lower() for cat in other_categories]:
                        colorized_filename += f"{Fore.YELLOW}{word}{Style.RESET_ALL}"
                    else:
                        colorized_filename += word
                print(f"- {colorized_filename.strip()}")
                
                print("\n" + "-"*50)

                all_summaries.append({
                    'file_path': os.path.abspath(file_path),
                    'category': top_category,
                    'summary': summary_text_for_file
                })

                if file_action != 'none' and top_category in file_management_settings and file_management_settings[top_category].get('destination'):
                    settings = file_management_settings[top_category]
                    prefix = settings['prefix']
                    destination_folder = settings['destination']
                    
                    os.makedirs(destination_folder, exist_ok=True)

                    original_file_name_no_ext = os.path.basename(os.path.splitext(file_path)[0])
                    
                    if prefix:
                        new_file_name = f"{prefix}_{original_file_name_no_ext}"
                    else:
                        new_file_name = original_file_name_no_ext

                    destination_path = os.path.join(destination_folder, new_file_name)
                    
                    try:
                        destination_path_with_ext = None
                        if file_action == 'shortcut':
                            create_shortcut(file_path, destination_path)
                            print(f"Created shortcut for: {os.path.basename(file_path)} -> {destination_folder}")
                        else:
                            destination_path_with_ext = destination_path + os.path.splitext(file_path)[1]
                            if file_action == 'move':
                                if os.path.exists(file_path):
                                    shutil.move(file_path, destination_path_with_ext)
                                    print(f"Moved and renamed: {file_path} -> {destination_path_with_ext}")
                            elif file_action == 'copy':
                                if os.path.exists(file_path):
                                    shutil.copy2(file_path, destination_path_with_ext)
                                    print(f"Copied and renamed: {file_path} -> {destination_path_with_ext}")
                        
                        # NEW: Tag the destination file if it was moved or copied
                        if destination_path_with_ext and tagging_enabled and top_score * 100 >= tagging_threshold:
                            tag_file_properties(destination_path_with_ext, top_category)

                    except Exception as e:
                        print(f"Error during file operation '{file_action}' for {os.path.basename(file_path)}: {e}")
            else:
                print("Step 9.1: Skipping file - summarization failed.")
        else:
            print("Step 9.1: Skipping file - no text found in the specified range.")
    
    print("\n--- Step 10: All supported files processed ---")
    return all_summaries


# --- Main Execution ---
if __name__ == "__main__":
    try:
        CATEGORIES = load_and_edit_categories()
        
        settings_file = "settings.json"
        
        defaults = {
            'pdf_chunk_length': 500,
            'generic_chunk_length': 250,
            'max_summary_length': 150, 'min_summary_length': 40,
            'pdf_padding': 0, 'pdf_chunks': 1, 'generic_padding': 0, 'generic_chunks': 1,
            'confidence_threshold': 0, 'tagging_enabled_choice': 'n', 'tagging_threshold': 90,
            'add_prefix_choice': 'n', 'global_prefix_choice': 'n', 'global_prefix': '', 'special_char': '',
            'selected_file_types': ['.pdf', '.docx', '.txt', '.xlsx', '.pptx'],
            'file_action_choice': 'none',
            'last_list_file': '',
            'last_destination_folder': '',
            'source_choice': 'f'
        }

        if os.path.exists(settings_file):
            load_choice = input("Found saved settings. Load them as defaults for the following prompts? (y/n) [default: y]: ").strip().lower()
            if load_choice in ['y', 'yes', '']:
                try:
                    with open(settings_file, 'r') as f:
                        defaults.update(json.load(f))
                    print("Settings loaded successfully.")
                except Exception as e:
                    print(f"Error loading settings: {e}. Using hardcoded defaults.")
        
        # --- File Type Selection ---
        supported_types_str = ", ".join(['.pdf', '.docx', '.txt', '.xlsx', '.pptx'])
        default_types_str = ", ".join(defaults['selected_file_types'])
        while True:
            print(f"\nSupported file types: {supported_types_str}")
            types_input = input(f"Enter file types to process, separated by commas (default: {default_types_str}): ").strip().lower()
            if not types_input:
                selected_file_types = defaults['selected_file_types']
                break
            else:
                selected_file_types = [t.strip() for t in types_input.split(',') if t.strip()]
                if all(t in supported_types_str for t in selected_file_types):
                    break
                else:
                    print("Invalid file type entered. Please choose from the supported types.")
        
        # --- Source Selection ---
        source_prompt = f"Do you want to process files from a folder (f) or a list from a text file (l)? [default: {defaults['source_choice']}]: "
        source_choice_input = input(source_prompt).strip().lower()
        source_choice = source_choice_input if source_choice_input else defaults['source_choice']
        
        initial_file_list = []
        output_folder_path = ""
        list_file_path = defaults.get('last_list_file', '')
        last_destination_folder = defaults.get('last_destination_folder', '')

        if source_choice == 'f':
            while True:
                folder_path = input("Enter the folder path to scan: ").strip()
                if os.path.isdir(folder_path):
                    break
                print("\nError: The provided path is not a valid directory. Please try again.")

            # --- Destination Choice for Folder Scans ---
            while True:
                dest_choice_prompt = "Where should organized files be sent? [S]ubdirectory of the current location, or [D]efined location?: "
                dest_choice = input(dest_choice_prompt).strip().lower()
                if dest_choice in ['s', 'd']:
                    break
                print("Invalid choice. Please enter 'S' or 'D'.")
            
            if dest_choice == 's':
                output_folder_path = folder_path
            else: # dest_choice == 'd'
                while True:
                    prompt = f"Enter the destination folder path for all output [default: {last_destination_folder}]: "
                    defined_path_input = input(prompt).strip()
                    defined_path = defined_path_input if defined_path_input else last_destination_folder
                    if defined_path:
                        output_folder_path = defined_path
                        if not os.path.isdir(output_folder_path):
                            print(f"Destination folder '{output_folder_path}' does not exist. Creating it.")
                            os.makedirs(output_folder_path, exist_ok=True)
                        break
                    print("Path cannot be empty. Please enter a valid path.")

            scan_sub_input = input("Scan subdirectories of the source folder? (y/n) [default: n]: ").strip().lower()
            scan_subdirectories = scan_sub_input in ['y', 'yes']

            print("\nScanning for files... This may take a moment for large directories.")
            if scan_subdirectories:
                for root, _, files in os.walk(folder_path):
                    for file in files:
                        if os.path.splitext(file)[1].lower() in selected_file_types:
                            initial_file_list.append(os.path.join(root, file))
            else:
                for file in os.listdir(folder_path):
                    file_path_check = os.path.join(folder_path, file)
                    if os.path.isfile(file_path_check) and os.path.splitext(file)[1].lower() in selected_file_types:
                        initial_file_list.append(file_path_check)

        elif source_choice == 'l':
            while True:
                prompt = f"Enter the path to the text file with the list of files [default: {list_file_path}]: "
                list_file_path_input = input(prompt).strip()
                list_file_path = list_file_path_input if list_file_path_input else list_file_path
                
                if os.path.exists(list_file_path) and os.path.isfile(list_file_path):
                    while True:
                        prompt = f"Enter a destination folder path for summaries and categorized files [default: {last_destination_folder}]: "
                        defined_path_input = input(prompt).strip()
                        defined_path = defined_path_input if defined_path_input else last_destination_folder
                        if defined_path:
                            output_folder_path = defined_path
                            if not os.path.isdir(output_folder_path):
                                print(f"Destination folder '{output_folder_path}' does not exist. Creating it.")
                                os.makedirs(output_folder_path, exist_ok=True)
                            break
                        print("Path cannot be empty. Please enter a valid path.")
                    break
                else:
                    print("\nError: The provided path is not a valid file. Please try again.")
                    list_file_path = ''

            print(f"\nScanning for files from the list file: {list_file_path}")
            with open(list_file_path, 'r', encoding='utf-8') as f:
                for line in f:
                    file_path = line.strip()
                    if file_path and os.path.exists(file_path) and os.path.isfile(file_path):
                        if os.path.splitext(file_path)[1].lower() in selected_file_types:
                            initial_file_list.append(file_path)
                        else:
                            print(f"Skipping unsupported file type: {file_path}")
                    else:
                        print(f"Warning: Could not find file specified in list: {file_path}")
        else:
            print("Invalid choice. Exiting.")
            sys.exit()

        # --- File Action Choice ---
        action_map = {'m': 'move', 'c': 'copy', 's': 'shortcut', 'd': 'none'}
        default_action = defaults['file_action_choice']
        while True:
            action_prompt = f"Do you want to [M]ove, [C]opy, create [S]hortcut, or [D]o nothing? [default: {default_action}]: "
            action_input = input(action_prompt).strip().lower()
            if not action_input:
                file_action_to_use = default_action
                break
            elif action_input in action_map:
                file_action_to_use = action_map[action_input]
                break
            else:
                print("Invalid choice. Please enter M, C, S, or D.")

        print("\n--- Configuration Settings ---")
        # --- Independent Chunk Lengths ---
        while True:
            try:
                chunk_len_input = input(f"Enter non-PDF chunk length (words per chunk, default: {defaults['generic_chunk_length']}): ").strip()
                generic_chunk_length_to_use = int(chunk_len_input) if chunk_len_input else defaults['generic_chunk_length']
                if generic_chunk_length_to_use > 0:
                    break
                else:
                    print("Please enter a number greater than 0.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")
        
        while True:
            try:
                chunk_len_input = input(f"Enter PDF chunk length (words per chunk, default: {defaults['pdf_chunk_length']}): ").strip()
                pdf_chunk_length_to_use = int(chunk_len_input) if chunk_len_input else defaults['pdf_chunk_length']
                if pdf_chunk_length_to_use > 0:
                    break
                else:
                    print("Please enter a number greater than 0.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")

        while True:
            try:
                max_sum_input = input(f"Enter max summary length (words, default: {defaults['max_summary_length']}): ").strip()
                max_summary_length_to_use = int(max_sum_input) if max_sum_input else defaults['max_summary_length']
                if max_summary_length_to_use > 0:
                    break
                else:
                    print("Please enter a number greater than 0.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")

        while True:
            try:
                min_sum_input = input(f"Enter min summary length (words, default: {defaults['min_summary_length']}): ").strip()
                min_summary_length_to_use = int(min_sum_input) if min_sum_input else defaults['min_summary_length']
                if min_summary_length_to_use > 0 and min_summary_length_to_use <= max_summary_length_to_use:
                    break
                else:
                    print(f"Please enter a number greater than 0 and no more than {max_summary_length_to_use}.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")

        print(f"\n--- PDF Settings ---")
        while True:
            try:
                padding_input = input(f"Enter number of initial chunks to skip (padding, default: {defaults['pdf_padding']}): ").strip()
                pdf_padding_to_use = int(padding_input) if padding_input else defaults['pdf_padding']
                if pdf_padding_to_use >= 0:
                    break
                else:
                    print("Please enter a non-negative number.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")
        
        while True:
            try:
                chunks_input = input(f"Enter number of chunks to process after padding (default: {defaults['pdf_chunks']}): ").strip()
                pdf_chunks_to_use = int(chunks_input) if chunks_input else defaults['pdf_chunks']
                if pdf_chunks_to_use > 0:
                    break
                else:
                    print("Please enter a number greater than 0.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")

        print(f"\n--- Settings for Other File Types ---")
        while True:
            try:
                padding_input = input(f"Enter number of initial chunks to skip (padding, default: {defaults['generic_padding']}): ").strip()
                generic_padding_to_use = int(padding_input) if padding_input else defaults['generic_padding']
                if generic_padding_to_use >= 0:
                    break
                else:
                    print("Please enter a non-negative number.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")

        while True:
            try:
                chunks_input = input(f"Enter number of chunks to process after padding (default: {defaults['generic_chunks']}): ").strip()
                generic_chunks_to_use = int(chunks_input) if chunks_input else defaults['generic_chunks']
                if generic_chunks_to_use > 0:
                    break
                else:
                    print("Please enter a number greater than 0.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")
        
        while True:
            try:
                threshold_input = input(f"\nEnter confidence threshold for categorization (0-100, default: {defaults['confidence_threshold']}): ").strip()
                confidence_threshold_to_use = int(threshold_input) if threshold_input else defaults['confidence_threshold']
                if 0 <= confidence_threshold_to_use <= 100:
                    break
                else:
                    print("Please enter a number between 0 and 100.")
            except ValueError:
                print("Invalid input. Please enter a valid number.")
        
        tagging_enabled_choice = input(f"\nTag files with category name in file properties? (y/n) [default: {defaults['tagging_enabled_choice']}]): 

").strip().lower() or defaults['tagging_enabled_choice']
        tagging_threshold_to_use = defaults['tagging_threshold']
        if tagging_enabled_choice in ['y', 'yes']:
            while True:
                try:
                    tag_thresh_input = input(f"Enter confidence threshold for tagging (0-100, default: {defaults['tagging_threshold']}): ").strip()
                    tagging_threshold_to_use = int(tag_thresh_input) if tag_thresh_input else defaults['tagging_threshold']
                    if 0 <= tagging_threshold_to_use <= 100:
                        break
                    else:
                        print("Please enter a number between 0 and 100.")
                except ValueError:
                    print("Invalid input. Please enter a valid number.")
        
        # Prefix settings only matter if a file action is chosen
        add_prefix_choice = defaults['add_prefix_choice']
        global_prefix_choice = defaults['global_prefix_choice']
        global_prefix = defaults['global_prefix']
        special_char = defaults['special_char']
        if file_action_to_use != 'none':
            add_prefix_choice = input(f"Add a prefix to organized files? (y/n) [default: {defaults['add_prefix_choice']}]): ").strip().lower() or 

defaults['add_prefix_choice']
            if add_prefix_choice in ['y', 'yes']:
                global_prefix_choice = input(f"Apply a global prefix to all? (y/n) [default: {defaults['global_prefix_choice']}]): ").strip().lower() or 

defaults['global_prefix_choice']
                if global_prefix_choice in ['y', 'yes']:
                    while True:
                        global_prefix = input(f"Enter a prefix (max 12 characters, default: {defaults['global_prefix']}): ").strip() or defaults

['global_prefix']
                        if len(global_prefix) <= 12:
                            break
                        else:
                            print(f"The prefix must be 12 characters or less. Please try again.")
                    
                    while True:
                        special_char = input(f"Enter a special character to follow the prefix (1 to 3 characters, default: {defaults['special_char']}): 

").strip() or defaults['special_char']
                        if 1 <= len(special_char) <= 3:
                            break
                        else:
                            print("Please enter 1 to 3 special characters.")

        settings_to_save = {
            'pdf_chunk_length': pdf_chunk_length_to_use,
            'generic_chunk_length': generic_chunk_length_to_use,
            'max_summary_length': max_summary_length_to_use,
            'min_summary_length': min_summary_length_to_use,
            'pdf_padding': pdf_padding_to_use,
            'pdf_chunks': pdf_chunks_to_use,
            'generic_padding': generic_padding_to_use,
            'generic_chunks': generic_chunks_to_use,
            'confidence_threshold': confidence_threshold_to_use,
            'tagging_enabled_choice': tagging_enabled_choice,
            'tagging_threshold': tagging_threshold_to_use,
            'add_prefix_choice': add_prefix_choice,
            'global_prefix_choice': global_prefix_choice,
            'global_prefix': global_prefix,
            'special_char': special_char,
            'selected_file_types': selected_file_types,
            'file_action_choice': file_action_to_use,
            'last_list_file': list_file_path if source_choice == 'l' else defaults['last_list_file'],
            'last_destination_folder': output_folder_path,
            'source_choice': source_choice
        }
        with open(settings_file, 'w') as f:
            json.dump(settings_to_save, f, indent=4)
        print("\nSettings saved to settings.json for the next session.")

        file_management_settings = {}
        if file_action_to_use != 'none' and output_folder_path:
            if add_prefix_choice in ['y', 'yes']:
                if global_prefix_choice in ['y', 'yes']:
                    for category in CATEGORIES:
                        if category != "Other":
                            file_management_settings[category] = {
                                'prefix': f"{global_prefix}{special_char}",
                                'destination': os.path.join(output_folder_path, category)
                            }
                else:
                    print("Using default prefixes for organized files.")
                    for category in CATEGORIES:
                        if category != "Other":
                            file_management_settings[category] = {
                                'prefix': category.replace(' ', '')[:12].upper(),
                                'destination': os.path.join(output_folder_path, category)
                            }
            else:
                print("Files will be organized without a prefix.")
                for category in CATEGORIES:
                    if category != "Other":
                        file_management_settings[category] = {
                            'prefix': '',
                            'destination': os.path.join(output_folder_path, category)
                        }
        
        files_to_process = initial_file_list

        print("\nStep 11: Folder path is valid. Starting the batch process...")
        
        all_summaries = process_files_in_folder(
            files_to_process, CATEGORIES, 
            pdf_padding_to_use, pdf_chunks_to_use, 
            generic_padding_to_use, generic_chunks_to_use, 
            file_management_settings, confidence_threshold_to_use, 
            pdf_chunk_length_to_use, generic_chunk_length_to_use, 
            max_summary_length_to_use, min_summary_length_to_use, 
            tagging_enabled_choice in ['y', 'yes'], tagging_threshold_to_use, 
            file_action_to_use, output_folder_path
        )

        if all_summaries and output_folder_path:
            print("\n--- Saving all summaries to a single file ---")
            consolidated_summary_file = os.path.join(output_folder_path, "all_summaries.txt")
            with open(consolidated_summary_file, "w", encoding="utf-8") as f:
                for summary_data in all_summaries:
                    f.write("="*50 + "\n")
                    f.write(f"File: {summary_data['file_path']}\n")
                    f.write(f"Category: {summary_data['category']}\n\n")
                    f.write(f"Summary:\n{summary_data['summary']}\n")
            print(f"All summaries consolidated into: {consolidated_summary_file}")
            
    except KeyboardInterrupt:
        print("\nProcess interrupted by user. Exiting gracefully.")
