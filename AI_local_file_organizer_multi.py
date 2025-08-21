# AI_local_file_organizer_multi.py.py
import os
import sys
import psutil # For system memory check
import fitz  # PyMuPDF for PDF
from docx import Document
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook
import torch
from transformers import pipeline
import shutil
import json
import colorama
from colorama import Fore, Back, Style
import re
import multiprocessing
import threading
from tqdm import tqdm # For the progress bar
import time
import uuid

# Suppress TensorFlow/oneDNN warnings
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '2'
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'

# Initialize colorama to auto-reset styles after each print
colorama.init(autoreset=True)

# --- Global Variables for Models ---
classifier = None # Summarizer is loaded in workers
settings = {}
# Constants from old script for console output
MAX_OUTPUT_LINES = 5
MAX_LINE_WIDTH = 120

# --- Helper Functions ---
def sanitize_for_path(name):
    sanitized = re.sub(r'[<>:"/\\|?*]', '', name)
    return sanitized[:50].strip()

# --- MODIFIED FUNCTION ---
def create_shortcut(source_path, shortcut_path_no_ext, print_lock):
    """Creates a shortcut, now with verbose output if it already exists."""
    try:
        if sys.platform == "win32":
            import win32com.client
            shortcut_path = shortcut_path_no_ext + ".lnk"
            if os.path.exists(shortcut_path):
                with print_lock:
                    tqdm.write(f"{Fore.YELLOW}Action:{Style.RESET_ALL} Shortcut for '{os.path.basename(source_path)}' already exists. Skipping.")
                return # Exit verbosely
            shell = win32com.client.Dispatch("WScript.Shell")
            shortcut = shell.CreateShortCut(shortcut_path)
            shortcut.TargetPath = os.path.abspath(source_path)
            shortcut.save()
        else: # for macOS/Linux
            if os.path.lexists(shortcut_path_no_ext):
                with print_lock:
                    tqdm.write(f"{Fore.YELLOW}Action:{Style.RESET_ALL} Symlink for '{os.path.basename(source_path)}' already exists. Skipping.")
                return
            os.symlink(os.path.abspath(source_path), shortcut_path_no_ext)
    except Exception as e:
        with print_lock:
            tqdm.write(f"{Fore.RED}Failed to create shortcut for '{os.path.basename(source_path)}': {e}{Style.RESET_ALL}")


def tag_file_properties(file_path, primary_category, secondary_categories, print_lock):
    """Adds the determined categories to the file's metadata keywords/tags."""
    with print_lock:
        try:
            if not os.access(file_path, os.W_OK):
                tqdm.write(f"{Fore.RED}Tagging skipped: No write permission for '{os.path.basename(file_path)}'.{Style.RESET_ALL}")
                return
            
            ext = os.path.splitext(file_path)[1].lower()
            keyword_parts = []
            if primary_category and primary_category != "Other":
                keyword_parts.append(f"Primary: {primary_category}")
            for cat in secondary_categories:
                keyword_parts.append(f"Secondary: {cat}")
            
            keyword_string = "; ".join(keyword_parts)
            if not keyword_string: return

            if ext == ".pdf":
                doc = fitz.open(file_path)
                doc.set_metadata({"keywords": keyword_string})
                doc.save(file_path, incremental=True, encryption=fitz.PDF_ENCRYPT_KEEP)
                doc.close()
            elif ext in [".docx", ".pptx", ".xlsx"]:
                obj = None
                if ext == ".docx": obj = Document(file_path)
                elif ext == ".pptx": obj = Presentation(file_path)
                elif ext == ".xlsx": obj = load_workbook(file_path)
                
                if ext != ".xlsx": obj.core_properties.keywords = keyword_string
                else: obj.properties.keywords = keyword_string
                obj.save(file_path)
            else: return
            
            tqdm.write(f"{Fore.CYAN}Tagged '{os.path.basename(file_path)}' with: {keyword_string}{Style.RESET_ALL}")

        except Exception as e:
            tqdm.write(f"{Fore.RED}Tagging Error for {os.path.basename(file_path)}: {e}{Style.RESET_ALL}")


def load_and_edit_categories():
    category_file = "categories.json"
    default_categories = ["a document related to cancer or medical oncology.", "a document related to modern global politics, international relations, or war.", "a document related to public health, epidemiology, or disease trends.", "a document related to neurology, including the brain or spine.", "a document related to cardiology or the heart.", "information about herbal remedies, vitamins, or dietary supplements.", "a document about banking, stock market, business, finance, or economics.", "a document about high-tech, software, or information technology.", "a document about natural sciences like geology, astronomy, or climate science.", "a document related to archaeology, or human evolution.", "a document related to ancient history, religion, or mythology."]
    categories = default_categories
    try:
        if os.path.exists(category_file):
            with open(category_file, 'r', encoding='utf-8') as f: categories = json.load(f)
            print(f"{Fore.CYAN}Step 1.2: Loaded categories from previous session.{Style.RESET_ALL}")
    except Exception as e:
        print(f"{Fore.LIGHTBLACK_EX}WARNING: Error loading categories file: {e}. Using default categories.{Style.RESET_ALL}")

    print(f"\n{Fore.LIGHTYELLOW_EX}--- Current Categories ---{Style.RESET_ALL}")
    for i, cat in enumerate(categories): print(f"{Fore.LIGHTYELLOW_EX}[{i+1}] {cat}{Style.RESET_ALL}")
    print(f"{Fore.LIGHTYELLOW_EX}--------------------------{Style.RESET_ALL}")

    if input(f"{Fore.GREEN}Do you want to edit these categories? (y/n) [default: no]: {Style.RESET_ALL}").strip().lower() in ['y', 'yes']:
        instructions = f"\n{Fore.LIGHTBLACK_EX}Editing categories: 'add <name>', 'remove <num>', 'edit <num> <name>', 'list', 'done'{Style.RESET_ALL}"
        print(instructions)
        while True:
            command = input("> ").strip()
            parts = command.split()
            if not parts: continue
            action = parts[0].lower()
            if action == 'done': break
            elif action == 'add' and len(parts) >= 2:
                categories.append(" ".join(parts[1:])); print(f"{Fore.CYAN}Added: {' '.join(parts[1:])}{Style.RESET_ALL}")
            elif action == 'remove' and len(parts) == 2 and parts[1].isdigit():
                idx = int(parts[1]) - 1
                if 0 <= idx < len(categories): print(f"{Fore.RED}Removed: {categories.pop(idx)}{Style.RESET_ALL}")
                else: print(f"{Fore.RED}Invalid number.{Style.RESET_ALL}")
            elif action == 'edit' and len(parts) >= 3 and parts[1].isdigit():
                idx = int(parts[1]) - 1
                if 0 <= idx < len(categories):
                    old_cat = categories[idx]
                    categories[idx] = " ".join(parts[2:])
                    print(f"{Fore.CYAN}Edited '{old_cat}' to '{categories[idx]}'{Style.RESET_ALL}")
                else: print(f"{Fore.RED}Invalid number.{Style.RESET_ALL}")
            elif action == 'list':
                for i, cat in enumerate(categories): print(f"{Fore.LIGHTWHITE_EX}[{i+1}] {cat}{Style.RESET_ALL}")
            else: print(f"{Fore.RED}Invalid command. {instructions}{Style.RESET_ALL}")

    with open(category_file, 'w', encoding='utf-8') as f: json.dump(categories, f, indent=4)
    print(f"\n{Fore.CYAN}Step 1.3: Categories saved.{Style.RESET_ALL}")
    return categories

def initialize_models():
    """Initializes the classifier model in the main process."""
    global classifier
    print(f"\n{Fore.CYAN}Step 1: Loading AI models...{Style.RESET_ALL}")
    device = "cuda" if torch.cuda.is_available() else "cpu"
    print(f"{Fore.CYAN}Device set to use {'GPU' if device == 'cuda' else 'CPU'}.{Style.RESET_ALL}")
    classifier_threads = settings.get('classifier_threads', 3)
    if device == "cpu":
        torch.set_num_threads(classifier_threads)
        print(f"{Fore.CYAN}Classifier set to use {classifier_threads} threads for its tasks.{Style.RESET_ALL}")
    try:
        classifier = pipeline("zero-shot-classification", model="MoritzLaurer/xtremedistil-l6-h256-zeroshot-v1.1-all-33", device=device, multi_label=True)
        print(f"{Fore.CYAN}Step 1.1: Classification model loaded.{Style.RESET_ALL}")
    except Exception as e:
        print(f"{Fore.LIGHTBLACK_EX}WARNING: Error loading classification model: {e}{Style.RESET_ALL}"); classifier = None

def read_entire_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            with fitz.open(file_path) as doc: text = "".join(page.get_text() for page in doc)
        elif ext == ".docx":
            doc = Document(file_path)
            text = "\n".join(para.text for para in doc.paragraphs)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f: text = f.read()
        elif ext == ".xlsx":
            df = pd.read_excel(file_path, engine='openpyxl')
            text = df.to_string()
        elif ext == ".pptx":
            prs = Presentation(file_path)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    except Exception:
        return "__READ_ERROR__"
    return text

# --- CONSOLE DISPLAY MANAGER ---

def display_final_summary(result_data, classification_results, print_lock):
    """
    Re-creates the detailed console output from your original script using tqdm.write().
    """
    with print_lock:
        file_path = result_data['source_path']
        raw_text = result_data['raw_text']
        culled_text = result_data['culled_text']
        summary_text = result_data['summary']
        
        if settings.get('show_chunk_verbiage', 'n').lower() == 'y':
            tqdm.write(f"\n{Fore.MAGENTA}--- Raw Text Verbiage ({len(raw_text.split())} words) ---{Style.RESET_ALL}")
            for line in raw_text.splitlines()[:MAX_OUTPUT_LINES]:
                tqdm.write(f"{Fore.MAGENTA}{line[:MAX_LINE_WIDTH]}")
            tqdm.write(f"{Fore.MAGENTA}--------------------------------------------------{Style.RESET_ALL}")
            
            tqdm.write(f"\n{Fore.CYAN}--- Culled Text for Summarizer ({len(culled_text.split())} words) ---{Style.RESET_ALL}")
            for line in culled_text.splitlines()[:MAX_OUTPUT_LINES]:
                tqdm.write(f"{Fore.CYAN}{line[:MAX_LINE_WIDTH]}")
            tqdm.write(f"{Fore.CYAN}--------------------------------------------------{Style.RESET_ALL}")

        tqdm.write(f"\n{Fore.GREEN}--- Plain Text Summary ---{Style.RESET_ALL}")
        for line in summary_text.splitlines()[:MAX_OUTPUT_LINES]:
            tqdm.write(f"{Fore.LIGHTGREEN_EX}- {line[:MAX_LINE_WIDTH]}")
        tqdm.write(f"{Fore.GREEN}--------------------------{Style.RESET_ALL}")

        primary_category = classification_results['primary_category']
        secondary_categories = classification_results['secondary_categories']
        original_primary_category = classification_results['original_primary_category']
        original_primary_score = classification_results['original_primary_score']
        
        tqdm.write(f"\n{Fore.GREEN}File: {Fore.LIGHTWHITE_EX}{file_path}{Style.RESET_ALL}")
        category_color = Fore.YELLOW if primary_category == "Other" else Fore.GREEN
        if primary_category == "Other" and original_primary_category != "Other":
            tqdm.write(f"{Fore.GREEN}Primary Category:{Style.RESET_ALL} {category_color}{Style.BRIGHT}{primary_category} ({Fore.LIGHTBLACK_EX}was '{original_primary_category}' at {original_primary_score*100:.2f}%{Style.RESET_ALL}){Style.RESET_ALL}")
        else:
            tqdm.write(f"{Fore.GREEN}Primary Category:{Style.RESET_ALL} {category_color}{Style.BRIGHT}{primary_category} ({Fore.LIGHTWHITE_EX}{classification_results['primary_score']*100:.2f}%{Style.RESET_ALL})")
        
        if secondary_categories:
            colors = [Fore.CYAN, Fore.LIGHTGREEN_EX]
            colored_cats = [f"{colors[i % 2]}{cat} ({Fore.YELLOW}{score*100:.2f}%)" for i, (cat, score) in enumerate(secondary_categories)]
            tqdm.write(f"{Fore.GREEN}Secondary Categories:{Style.RESET_ALL} {', '.join(colored_cats)}")

        tqdm.write("\n" + "="*60)


# --- FILE-BASED WORKER LOGIC ---

def preprocessor_thread(source_files, job_dirs, num_workers, settings_data, print_lock, stop_event):
    for source_path in source_files:
        if stop_event.is_set(): break
        while True:
            pending_summaries = len(os.listdir(job_dirs['summarize']))
            pending_classifications = len(os.listdir(job_dirs['classify']))
            if (pending_summaries + pending_classifications) < (num_workers * 2):
                break
            time.sleep(0.5)

        raw_text = read_entire_file(source_path)
        if not raw_text or raw_text == "__READ_ERROR__":
            with print_lock:
                tqdm.write(f"{Fore.YELLOW}Skipped (unreadable):{Style.RESET_ALL} {os.path.basename(source_path)}")
            continue
            
        culled_text = raw_text # In a future step, your detailed culling logic would go here
        
        job_content = json.dumps({'source_path': source_path, 'raw_text': raw_text, 'culled_text': culled_text})
        job_filepath = os.path.join(job_dirs['summarize'], f"{uuid.uuid4()}.job")
        with open(job_filepath, 'w', encoding='utf-8') as f: f.write(job_content)
        
        output_size = os.path.getsize(job_filepath)
        with print_lock:
            tqdm.write(f"{Fore.BLUE}Preprocessor:{Style.RESET_ALL} Created job for '{os.path.basename(source_path)}' ({output_size} bytes)")

def summarizer_process(job_dirs, settings_data, print_lock, stop_event):
    device = "cuda" if torch.cuda.is_available() else "cpu"
    summarizer_pipeline = pipeline("summarization", model="t5-small", device=device)
    pid = os.getpid()

    while not stop_event.is_set():
        job_file = None
        try:
            for fname in os.listdir(job_dirs['summarize']):
                if fname.endswith('.job'):
                    base_path = os.path.join(job_dirs['summarize'], fname)
                    locked_path = base_path + '.locked'
                    try:
                        os.rename(base_path, locked_path)
                        job_file = locked_path
                        break
                    except (OSError, PermissionError): continue
            
            if not job_file:
                time.sleep(0.2)
                continue

            with open(job_file, 'r', encoding='utf-8') as f: job_data = json.load(f)
            
            text_to_summarize = job_data['culled_text']
            summary_result = summarizer_pipeline(text_to_summarize, max_length=settings_data['max_summary_length'], min_length=settings_data['min_summary_length'], truncation=True)
            summary_text = summary_result[0]['summary_text']

            result_content = json.dumps({
                'source_path': job_data['source_path'], 
                'raw_text': job_data['raw_text'],
                'culled_text': text_to_summarize,
                'summary': summary_text
            })
            result_filepath = os.path.join(job_dirs['classify'], os.path.basename(job_file).replace('.job.locked', '.sum'))
            with open(result_filepath, 'w', encoding='utf-8') as f: f.write(result_content)
            os.remove(job_file)

        except Exception as e:
            with print_lock: tqdm.write(f"{Fore.RED}Error in summarizer {pid}: {e}")
            if job_file and os.path.exists(job_file): os.remove(job_file)

def classifier_thread(total_jobs, job_dirs, categories, settings_data, print_lock, stop_event):
    processed_count = 0
    with tqdm(total=total_jobs, desc="Overall Progress") as pbar:
        while processed_count < total_jobs:
            if stop_event.is_set(): break
            job_file = None
            try:
                for fname in os.listdir(job_dirs['classify']):
                    if fname.endswith('.sum'):
                        base_path = os.path.join(job_dirs['classify'], fname)
                        locked_path = base_path + '.locked'
                        try:
                            os.rename(base_path, locked_path)
                            job_file = locked_path
                            break
                        except (OSError, PermissionError): continue
                
                if not job_file:
                    time.sleep(0.2)
                    continue

                with open(job_file, 'r', encoding='utf-8') as f: result_data = json.load(f)
                summary_text = result_data['summary']
                
                if classifier:
                    results = classifier(summary_text, candidate_labels=categories, multi_label=True)
                else:
                    results = {'labels': ['Other'], 'scores': [1.0]}

                classification_output = {}
                original_primary_category, original_primary_score = results['labels'][0], results['scores'][0]
                primary_category, secondary_categories = "Other", []
                if original_primary_score * 100 >= settings_data['confidence_threshold']:
                    primary_category = original_primary_category
                if settings_data.get('assign_secondary_categories', 'n').lower() == 'y':
                    if len(results['labels']) > 1 and results['scores'][1] * 100 >= settings_data['secondary_confidence_threshold']:
                        secondary_categories.append((results['labels'][1], results['scores'][1]))

                classification_output.update({
                    'primary_category': primary_category, 'primary_score': original_primary_score,
                    'secondary_categories': secondary_categories, 'original_primary_category': original_primary_category,
                    'original_primary_score': original_primary_score
                })
                
                display_final_summary(result_data, classification_output, print_lock)
                
                if settings_data.get('tagging_enabled_choice', 'n').lower() == 'y':
                    if primary_category and primary_category != "Other":
                        sec_cat_names = [cat for cat, score in secondary_categories]
                        tag_file_properties(result_data['source_path'], primary_category, sec_cat_names, print_lock)

                # --- MODIFIED: FILE ACTION LOGIC ---
                file_action = settings_data.get('file_action_choice', 'none')
                if file_action != 'none':
                    source_path = result_data['source_path']
                    output_folder = settings_data['output_folder_path']

                    if primary_category and primary_category != "Other":
                        dest_dir = os.path.join(output_folder, sanitize_for_path(primary_category))
                        os.makedirs(dest_dir, exist_ok=True)
                        
                        base_name = os.path.basename(source_path)
                        dest_path = os.path.join(dest_dir, base_name)

                        try:
                            if file_action == 'shortcut':
                                shortcut_path_no_ext = os.path.join(dest_dir, os.path.splitext(base_name)[0])
                                # Pass the print_lock to the function
                                create_shortcut(source_path, shortcut_path_no_ext, print_lock)
                                with print_lock:
                                    tqdm.write(f"{Fore.CYAN}Action:{Style.RESET_ALL} Created shortcut for '{base_name}' in '{primary_category}'")
                            elif file_action == 'copy':
                                if not os.path.exists(dest_path):
                                    shutil.copy2(source_path, dest_path)
                                    with print_lock:
                                        tqdm.write(f"{Fore.CYAN}Action:{Style.RESET_ALL} Copied '{base_name}' to '{primary_category}'")
                            elif file_action == 'move':
                                if os.path.exists(source_path):
                                    shutil.move(source_path, dest_path)
                                    with print_lock:
                                        tqdm.write(f"{Fore.CYAN}Action:{Style.RESET_ALL} Moved '{base_name}' to '{primary_category}'")
                        except Exception as e:
                            with print_lock:
                                tqdm.write(f"{Fore.RED}File Action Error:{Style.RESET_ALL} Could not {file_action} '{base_name}'. Reason: {e}")
                
                os.remove(job_file)
                processed_count += 1
                pbar.update(1)

            except Exception as e:
                with print_lock: tqdm.write(f"{Fore.RED}Error in classifier: {e}")
                if job_file and os.path.exists(job_file): os.remove(job_file)
    stop_event.set()

# --- Main Execution ---
if __name__ == "__main__":
    multiprocessing.freeze_support()
    settings_file = "settings.json"
    defaults = {'classifier_threads': 3, 'pdf_chunk_length': 2048, 'generic_chunk_length': 250, 'max_summary_length': 150, 'min_summary_length': 40, 'confidence_threshold': 50, 'secondary_confidence_threshold': 20, 'tagging_enabled_choice': 'n', 'add_prefix_choice': 'n', 'global_prefix_choice': 'n', 'global_prefix': '', 'special_char': '', 'selected_file_types': ['.pdf', '.docx', '.txt', '.xlsx', '.pptx'], 'file_action_choice': 'none', 'last_list_file': '', 'last_destination_folder': '', 'source_choice': 'f', 'preprocess_choice': 'n', 'summarizer_timeout': 60, 'long_pdf_threshold': 10, 'long_text_threshold': 1000, 'short_pdf_chunk_length': 1025, 'long_pdf_chunk_length': 4096, 'short_text_chunk_length': 1025, 'long_text_chunk_length': 4096, 'long_pdf_max_summary_length': 200, 'long_pdf_min_summary_length': 50, 'short_pdf_max_summary_length': 150, 'short_pdf_min_summary_length': 40, 'short_pdf_chunks_list': '1', 'long_pdf_chunks_list': '1', 'short_text_chunks_list': '1', 'long_text_chunks_list': '1', 'show_chunk_verbiage': 'n', 'filter_keywords': '', 'negative_filter_keywords': '', 'culling_regex': '', 'exclude_regex': '', 'culling_mode': 'c', 'user_threads': None, 'assign_secondary_categories': 'n'}
    settings.update(defaults)
    if os.path.exists(settings_file):
        print(f"{Fore.GREEN}Found saved settings. Load them? (y/n) [default: y]: {Style.RESET_ALL}", end="")
        if input().strip().lower() not in ['n', 'no']:
            try:
                with open(settings_file, 'r') as f: settings.update(json.load(f))
                print(f"{Fore.CYAN}Settings loaded successfully.{Style.RESET_ALL}")
            except Exception as e: print(f"{Fore.RED}Error loading settings: {e}.{Style.RESET_ALL}")
    try:
        CATEGORIES = load_and_edit_categories()
        print(f"{Fore.GREEN}Do you want to see the extended help content and regex examples? (y/n) [default: no]: {Style.RESET_ALL}", end="")
        if input().strip().lower() in ['y', 'yes']:
             print(f"\n{Fore.GREEN}# Advanced Regex Help: Common Patterns{Style.RESET_ALL}\n"
                   f"{Fore.CYAN}  {Fore.GREEN}Email:{Style.RESET_ALL} [a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\\.[a-zA-Z]{{2,}}\n"
                   f"{Fore.CYAN}  {Fore.GREEN}Phone (US):{Style.RESET_ALL} \\d{{3}}[-\\s]?\\d{{3}}[-\\s]?\\d{{4}}\n"
                   f"\n{Fore.GREEN}---{Style.RESET_ALL}")
        settings_to_save = settings.copy()
        print(f"{Fore.GREEN}Process files from a [f]older or a [l]ist file? [default: {settings['source_choice']}]: {Style.RESET_ALL}", end="")
        settings_to_save['source_choice'] = input().strip().lower() or settings['source_choice']
        initial_file_list, output_folder_path = [], ""
        last_destination_folder = settings['last_destination_folder']
        if settings_to_save['source_choice'] == 'f':
            print(f"{Fore.GREEN}Enter the folder path to scan: {Style.RESET_ALL}", end="")
            folder_path = input().strip() or '.'
            print(f"{Fore.GREEN}Send organized files to [s]ubdirectory or [d]efined location?: {Style.RESET_ALL}", end="")
            dest_choice = input().strip().lower()
            if dest_choice == 's': output_folder_path = os.path.join(folder_path, "Sorted")
            else:
                print(f"{Fore.GREEN}Enter destination folder [default: {last_destination_folder}]: {Style.RESET_ALL}", end="")
                output_folder_path = input().strip() or last_destination_folder
            print(f"{Fore.GREEN}Scan subdirectories? (y/n) [default: n]: {Style.RESET_ALL}", end="")
            if input().strip().lower() in ['y', 'yes']:
                for root, _, files in os.walk(folder_path):
                    for file in files: initial_file_list.append(os.path.join(root, file))
            else:
                for item in os.listdir(folder_path):
                    if os.path.isfile(os.path.join(folder_path, item)): initial_file_list.append(os.path.join(folder_path, item))
        elif settings_to_save['source_choice'] == 'l':
            print(f"{Fore.GREEN}Enter path to list file [default: {settings['last_list_file']}]: {Style.RESET_ALL}", end="")
            list_file_path = input().strip() or settings['last_list_file']
            settings_to_save['last_list_file'] = list_file_path
            print(f"{Fore.GREEN}Enter destination folder [default: {last_destination_folder}]: {Style.RESET_ALL}", end="")
            output_folder_path = input().strip() or last_destination_folder
            try:
                with open(list_file_path, 'r', encoding='utf-8') as f:
                    initial_file_list = [line.strip() for line in f if os.path.isfile(line.strip())]
            except FileNotFoundError: print(f"{Fore.RED}List file not found. Exiting.{Style.RESET_ALL}"); sys.exit(1)
        os.makedirs(output_folder_path, exist_ok=True)
        settings_to_save.update({'output_folder_path': output_folder_path, 'last_destination_folder': output_folder_path})
        
        action_map = {'m': 'move', 'c': 'copy', 's': 'shortcut', 'd': 'none'}
        default_action_key = next((k for k, v in action_map.items() if v == settings['file_action_choice']), 'd')
        print(f"{Fore.GREEN}Action: [M]ove, [C]opy, create [S]hortcut, or [D]o nothing? [default: {settings['file_action_choice']}]: {Style.RESET_ALL}", end="")
        settings_to_save['file_action_choice'] = action_map.get(input().strip().lower() or default_action_key, 'none')
        
        print(f"{Fore.GREEN}Pre-process by filename? (y/n) [default: {settings['preprocess_choice']}]: {Style.RESET_ALL}", end="")
        settings_to_save['preprocess_choice'] = input().strip().lower() or settings['preprocess_choice']

        print(f"\n{Fore.GREEN}--- Configuration Settings ---{Style.RESET_ALL}")
        print(f"\n{Fore.GREEN}# Core Settings{Style.RESET_ALL}")
        print(f"{Fore.CYAN}# These settings control the overall behavior of the summarizer and categorizer.{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Enter number of threads for classifier [default: {settings['classifier_threads']}]: {Style.RESET_ALL}", end="")
        settings_to_save['classifier_threads'] = int(input().strip() or settings['classifier_threads'])
        print(f"{Fore.GREEN}Enter summarizer timeout in seconds [default: {settings['summarizer_timeout']}]: {Style.RESET_ALL}", end="")
        settings_to_save['summarizer_timeout'] = int(input().strip() or settings['summarizer_timeout'])
        print(f"{Fore.GREEN}Enter primary confidence threshold (0-100) [default: {settings['confidence_threshold']}]: {Style.RESET_ALL}", end="")
        settings_to_save['confidence_threshold'] = int(input().strip() or settings['confidence_threshold'])
        print(f"{Fore.GREEN}Enter secondary confidence threshold (0-100) [default: {settings['secondary_confidence_threshold']}]: {Style.RESET_ALL}", end="")
        settings_to_save['secondary_confidence_threshold'] = int(input().strip() or settings['secondary_confidence_threshold'])
        print(f"{Fore.GREEN}Tag files with category keywords? (y/n) [default: {settings['tagging_enabled_choice']}]: {Style.RESET_ALL}", end="")
        settings_to_save['tagging_enabled_choice'] = input().strip().lower() or settings['tagging_enabled_choice']
        
        print(f"\n{Fore.GREEN}---{Style.RESET_ALL}")
        
        print(f"\n{Fore.GREEN}# Chunk & Text Settings{Style.RESET_ALL}")
        print(f"{Fore.CYAN}# These settings are used for parsing and preparing text from different file types.{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Enter non-PDF chunk length [default: {settings['generic_chunk_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['generic_chunk_length'] = int(input().strip() or settings['generic_chunk_length'])
        print(f"{Fore.GREEN}Enter SHORT PDF chunk length [default: {settings['short_pdf_chunk_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['short_pdf_chunk_length'] = int(input().strip() or settings['short_pdf_chunk_length'])
        print(f"{Fore.GREEN}Enter LONG PDF chunk length [default: {settings['long_pdf_chunk_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_chunk_length'] = int(input().strip() or settings['long_pdf_chunk_length'])
        print(f"{Fore.GREEN}Enter SHORT PDF max summary length [default: {settings['short_pdf_max_summary_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['short_pdf_max_summary_length'] = int(input().strip() or settings['short_pdf_max_summary_length'])
        print(f"{Fore.GREEN}Enter SHORT PDF min summary length [default: {settings['short_pdf_min_summary_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['short_pdf_min_summary_length'] = int(input().strip() or settings['short_pdf_min_summary_length'])
        print(f"{Fore.GREEN}Enter the LONG PDF threshold (in pages) [default: {settings['long_pdf_threshold']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_threshold'] = int(input().strip() or settings['long_pdf_threshold'])
        print(f"{Fore.GREEN}Enter LONG PDF max summary length [default: {settings['long_pdf_max_summary_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_max_summary_length'] = int(input().strip() or settings['long_pdf_max_summary_length'])
        print(f"{Fore.GREEN}Enter LONG PDF min summary length [default: {settings['long_pdf_min_summary_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_min_summary_length'] = int(input().strip() or settings['long_pdf_min_summary_length'])

        print(f"\n{Fore.GREEN}---{Style.RESET_ALL}")
        
        print(f"\n{Fore.GREEN}# Chunk Selection (comma-separated list/ranges){Style.RESET_ALL}")
        print(f"{Fore.CYAN}# These settings define which chunks of a document are read and processed.{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Enter chunks for SHORT PDFs (e.g., '1, 3-5') [default: {settings['short_pdf_chunks_list']}]: {Style.RESET_ALL}", end="")
        settings_to_save['short_pdf_chunks_list'] = input().strip() or settings['short_pdf_chunks_list']
        print(f"{Fore.GREEN}Enter chunks for LONG PDFs (e.g., '1-2, 10') [default: {settings['long_pdf_chunks_list']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_pdf_chunks_list'] = input().strip() or settings['long_pdf_chunks_list']
        print(f"{Fore.GREEN}Enter chunks for SHORT text files (e.g., '1') [default: {settings['short_text_chunks_list']}]: {Style.RESET_ALL}", end="")
        settings_to_save['short_text_chunks_list'] = input().strip() or settings['short_text_chunks_list']
        print(f"{Fore.GREEN}Enter chunks for LONG text files (e.g., '1-5') [default: {settings['long_text_chunks_list']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_text_chunks_list'] = input().strip() or settings['long_text_chunks_list']
        print(f"{Fore.GREEN}Enter the LONG TEXT threshold (in words) [default: {settings['long_text_threshold']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_text_threshold'] = int(input().strip() or settings['long_text_threshold'])
        print(f"{Fore.GREEN}Enter SHORT text chunk length [default: {settings['short_text_chunk_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['short_text_chunk_length'] = int(input().strip() or settings['short_text_chunk_length'])
        print(f"{Fore.GREEN}Enter LONG text chunk length [default: {settings['long_text_chunk_length']}]: {Style.RESET_ALL}", end="")
        settings_to_save['long_text_chunk_length'] = int(input().strip() or settings['long_text_chunk_length'])
        
        print(f"\n{Fore.GREEN}---{Style.RESET_ALL}")

        print(f"\n{Fore.GREEN}# Keyword & Regex Settings{Style.RESET_ALL}")
        print(f"{Fore.CYAN}# These settings provide fine-grained control over keyword filtering and text culling.{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Show extraction and summary text for each chunk? (y/n) [default: {settings['show_chunk_verbiage']}]: {Style.RESET_ALL}", end="")
        settings_to_save['show_chunk_verbiage'] = input().strip().lower() or settings['show_chunk_verbiage']
        
        print(f"\n{Fore.CYAN}# Culling Mode: Choose your method for pre-processing text for summarization.{Style.RESET_ALL}")
        print(f"{Fore.CYAN}#     'c' - Cull to complete sentences (Default).")
        print(f"{Fore.CYAN}#     's' - Extract the complete sentence(s) containing the keyword.")
        print(f"{Fore.CYAN}#     'b' - Extract the sentence before, the sentence containing, and the sentence after the keyword.")
        print(f"{Fore.CYAN}#     '2' - Extract two sentences before, the sentence containing, and two sentences after the keyword.")
        print(f"{Fore.CYAN}#     'p' - Extract the full paragraph containing the keyword.")
        print(f"{Fore.CYAN}#     'r' - Use a custom regex string to define the culling pattern.")
        print(f"{Fore.CYAN}#     'n' - No culling; summarize the entire extracted text.")
        print(f"{Fore.GREEN}Select a culling mode (c/s/b/2/p/r/n) [default: {settings['culling_mode']}]: {Style.RESET_ALL}", end="")
        settings_to_save['culling_mode'] = input().strip().lower() or settings['culling_mode']
        
        if settings_to_save['culling_mode'] in ['s', 'b', '2', 'p']:
            print(f"\n{Fore.CYAN}# These keywords are used to locate the relevant sentences/paragraphs for culling.{Style.RESET_ALL}")
            print(f"{Fore.GREEN}Enter filter keywords (comma-separated) [default: '{settings['filter_keywords']}']: {Style.RESET_ALL}", end="")
            settings_to_save['filter_keywords'] = input().strip() or settings['filter_keywords']
        
        print(f"\n{Fore.CYAN}# These keywords are used to skip processing files that are likely irrelevant.{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Enter keywords to filter OUT (comma-separated) [default: '{settings['negative_filter_keywords']}']: {Style.RESET_ALL}", end="")
        settings_to_save['negative_filter_keywords'] = input().strip() or settings['negative_filter_keywords']
        
        print(f"\n{Fore.CYAN}# This regex is used to exclude files from processing.{Style.RESET_ALL}")
        print(f"{Fore.GREEN}Enter regex to exclude files [default: '{settings['exclude_regex']}']: {Style.RESET_ALL}", end="")
        settings_to_save['exclude_regex'] = input().strip() or settings['exclude_regex']

        print(f"\n{Fore.GREEN}---{Style.RESET_ALL}")
        
        if settings_to_save['file_action_choice'] != 'none':
            print(f"{Fore.GREEN}Add a prefix to organized files? (y/n) [default: {settings['add_prefix_choice']}]: {Style.RESET_ALL}", end="")
            settings_to_save['add_prefix_choice'] = input().strip().lower() or settings['add_prefix_choice']
            if settings_to_save['add_prefix_choice'] in ['y', 'yes']:
                print(f"{Fore.GREEN}Apply a global prefix? (y/n) [default: {settings['global_prefix_choice']}]: {Style.RESET_ALL}", end="")
                settings_to_save['global_prefix_choice'] = input().strip().lower() or settings['global_prefix_choice']
                if settings_to_save['global_prefix_choice'] in ['y', 'yes']:
                    print(f"{Fore.GREEN}Enter a global prefix (max 12 chars) [default: {settings['global_prefix']}]: {Style.RESET_ALL}", end="")
                    settings_to_save['global_prefix'] = input().strip() or settings['global_prefix']
        
        num_cores = os.cpu_count()
        print(f"\n{Fore.GREEN}Your system has {num_cores} CPU cores. Enter the number of summarizer workers [-1 for max] [default: {settings.get('user_threads', 0)}]: {Style.RESET_ALL}", end="")
        user_threads_input = input().strip()
        settings_to_save['user_threads'] = int(user_threads_input or settings.get('user_threads') or 0)
        
        print(f"\n{Fore.GREEN}How many instances of the summarizer would you like to run? (1 instance is 1 core) [default: {settings_to_save.get('user_threads', 0)}]: {Style.RESET_ALL}", end="")
        user_threads_input = input().strip()
        settings_to_save['user_threads'] = int(user_threads_input or settings_to_save.get('user_threads') or 0)
        
        print(f"\n{Fore.GREEN}Assign secondary categories if confidence is met? (y/n) [default: {settings.get('assign_secondary_categories', 'n')}]: {Style.RESET_ALL}", end="")
        settings_to_save['assign_secondary_categories'] = input().strip().lower() or settings.get('assign_secondary_categories', 'n')
        
        settings.update(settings_to_save)
        with open(settings_file, 'w') as f: json.dump(settings, f, indent=4)
        print(f"\n{Fore.CYAN}Settings saved. Starting process...{Style.RESET_ALL}")
        
        # --- NEW PROCESSING ENGINE START ---
        initialize_models()
        files_to_process = [f for f in initial_file_list if os.path.splitext(f)[1].lower() in settings['selected_file_types']]
        total_jobs = len(files_to_process)
        if not files_to_process:
            print(f"{Fore.YELLOW}No files to process.{Style.RESET_ALL}")
            sys.exit(0)
            
        base_job_dir = os.path.join(output_folder_path, "_temp_jobs")
        job_dirs = {'summarize': os.path.join(base_job_dir, '2_summarize'), 'classify': os.path.join(base_job_dir, '3_classify')}
        if os.path.exists(base_job_dir): shutil.rmtree(base_job_dir)
        for d in job_dirs.values(): os.makedirs(d)

        num_workers = settings.get('user_threads', 0)
        if num_workers <= 0: num_workers = max(1, os.cpu_count() // 2 if num_workers == 0 else os.cpu_count())

        print_lock = multiprocessing.Lock()
        stop_event = multiprocessing.Event()

        preprocessor = threading.Thread(target=preprocessor_thread, args=(files_to_process, job_dirs, num_workers, settings, print_lock, stop_event))
        classifier_manager = threading.Thread(target=classifier_thread, args=(total_jobs, job_dirs, CATEGORIES, settings, print_lock, stop_event))
        summarizers = [multiprocessing.Process(target=summarizer_process, args=(job_dirs, settings, print_lock, stop_event)) for _ in range(num_workers)]

        print(f"{Fore.CYAN}Starting {num_workers} summarizers, 1 preprocessor, and 1 classifier...{Style.RESET_ALL}")
        preprocessor.start()
        classifier_manager.start()
        for p in summarizers: p.start()

        classifier_manager.join()
        preprocessor.join()
        for p in summarizers: p.join()

        if os.path.exists(base_job_dir): shutil.rmtree(base_job_dir)
        print(f"\n{Fore.GREEN}Processing complete.{Style.RESET_ALL}")

    except KeyboardInterrupt: 
        print(f"{Fore.LIGHTBLACK_EX}\nProcess interrupted by user.{Style.RESET_ALL}")
        stop_event.set()
    except Exception as e:
        import traceback
        print(f"{Fore.RED}\nAn unexpected error occurred: {e}{Style.RESET_ALL}")
        traceback.print_exc()
        time.sleep(15)
